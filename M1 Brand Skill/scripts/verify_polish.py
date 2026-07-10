#!/usr/bin/env python3
"""
Regression harness for the S1 cover and S6 quote polish rules.

Each case in CASES below builds a single-slide .pptx in testing/ and verifies
that every line fits within the placeholder width using actual Inter Variable Thin font
metrics. Add new cases by appending to CASES.

Usage:  python3 scripts/verify_polish.py
"""
import os
import shutil
import sys

BASE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(BASE)
sys.path.insert(0, BASE)

from pptx import Presentation
import pptx_helpers
from pptx_helpers import (
    duplicate_slide, delete_slide, set_cover_title, set_quote, set_section_break,
    set_content_bullets, set_content_title, set_content_slides,
    pick_content_slide_index, set_ph, set_slide_view, _text_width_pt,
    _get_ph_dims_from_xml, _find_section_support_shape,
)

TMPL    = os.path.join(ROOT, "assets", "M1-Presentation-Template_2026.pptx")
# Regression outputs go to a subfolder so they don't clutter the canonical
# demo deck (testing/M1-digital-private-bank.pptx). Created on first run.
OUT_DIR = os.path.join(ROOT, "testing", "_regression")
os.makedirs(OUT_DIR, exist_ok=True)

# S-template indices in the M1 2026 template (16 slides total)
COVER_SLIDE_INDEX         = 0   # S1
SECTION_LIGHT_SLIDE_INDEX = 2   # S3 — light mode
SECTION_DARK_SLIDE_INDEX  = 3   # S4 — dark mode
QUOTE_SLIDE_INDEX         = 5   # S6
N_TEMPLATE_SLIDES         = 18  # S1..S18 — all get deleted after our slide is appended

CASES = [
    {"kind": "cover", "out": "test-cover-01.pptx",
     "title": "M1 Intelligence: your always-on financial advisor"},
    {"kind": "cover", "out": "test-cover-02.pptx",
     "title": "Built for investors who think in decades"},
    {"kind": "section", "mode": "light", "out": "test-section-01-light-short.pptx",
     "title": "Why personal finance is broken"},
    {"kind": "section", "mode": "light", "out": "test-section-02-light-with-support.pptx",
     "title": "M1 is the third option",
     "supporting_text": "An integrated platform with intelligent automation — what private banks provide, for everyone."},
    {"kind": "section", "mode": "dark", "out": "test-section-03-dark-short.pptx",
     "title": "Platform traction"},
    {"kind": "section", "mode": "dark", "out": "test-section-04-dark-with-support.pptx",
     "title": "The M1 advantage",
     "supporting_text": "Where automation meets ownership, and discipline becomes the platform's job."},
    {"kind": "quote", "out": "test-quote-01.pptx",
     "quote": "Managing money used to feel overwhelming. With M1, it just runs.",
     "attribution": "M1 client"},
    {"kind": "quote", "out": "test-quote-02.pptx",
     "quote": "I tell M1 my strategy. It executes — without me, without emotion.",
     "attribution": "Jamie R., M1 client"},
    # Overflow case — deliberately enough content to force pagination across
    # two or more slides. Should produce a continuation slide titled
    # "<title> (continued)".
    {"kind": "content", "out": "test-content-04-overflow-paginates.pptx",
     "title": "Long-form deep dive that needs more than one slide of content",
     "bullet_mode": "default",
     "items": [
         (0, "First major theme covers cash flow fundamentals and why earning more on idle cash matters more than most investors realize"),
         (1, "High-yield cash beats brokerage sweep accounts by 3-5x for the same liquidity, and the gap compounds significantly over multi-year horizons even at modest balances."),
         (1, "Treasury exposure offers similar yields with sovereign-credit safety, ideal for emergency funds and short-term goals where principal protection matters more than upside."),
         (2, "M1 High-Yield Cash Account leads national averages"),
         (2, "Sweep into Treasuries with automated rebalancing"),
         (2, "Earn while you wait for investment opportunities"),
         (3, "FDIC insurance up to applicable limits"),
         (4, "Coverage extends across partner banks"),
         (0, "Second theme: balance sheet construction for long-term wealth-builders who think in decades not days"),
         (1, "Diversified portfolio construction is the single biggest predictor of outcome variance, dwarfing security selection and market timing combined over any 10+ year window."),
         (1, "M1 Pies let you express any allocation as a tree of weighted slices, with rebalancing handled automatically every time you contribute."),
         (2, "Build a Pie from scratch or clone an Expert Pie"),
         (2, "Dynamic Rebalancing keeps allocations on target"),
         (2, "Contribute as little or as much as you want"),
         (3, "Fractional shares mean every dollar invests"),
         (0, "Third theme: borrowing strategically without disrupting your long-term investment compounding"),
         (1, "Margin loans and personal loans give you access to liquidity without forced selling, preserving the compounding engine that drives long-term wealth."),
         (2, "M1 Margin Loan: rates competitive with institutional lenders"),
         (2, "M1 Personal Loan: fixed rates from 7.99% APR"),
         (3, "No origination fees, no late fees, no prepayment penalties"),
         (0, "Fourth theme: stewardship — protecting and optimizing what you've already built"),
         (1, "Trade windows and Dynamic Rebalancing combine to enforce discipline by design, protecting against the emotional decisions that erode returns during volatile markets."),
         (2, "Smart Transfers move money based on your rules"),
         (2, "Tax-loss harvesting integrated with portfolio strategy"),
         (3, "Aligns with your overall financial picture"),
     ]},
    # Short title — exercises the S7 (single-line title) code path.
    {"kind": "content", "out": "test-content-00-short-title-s7.pptx",
     "title": "Three pillars at a glance",
     "bullet_mode": "default",
     "items": [
         (0, "Three pillars of your financial life"),
         (1, "Each pillar is strong on its own. Together they create something none could deliver alone."),
         (2, "Cash flow: earn more on your money while staying liquid"),
         (2, "Balance sheet: invest for the long term, borrow when it makes sense"),
         (2, "Stewardship: protect what you've built and optimize how it grows"),
     ]},
    # Shared sample content for the three content-mode tests. Long title forces
    # S8 (so we also exercise the phrase-break + no-widow rule on the title).
    # Body content has multiple L1 sections and multiple consecutive L2 lines so
    # the L1↔L1 and L2↔L2 group-separator spacing is visible. Several lines are
    # written long enough to wrap to 2+ lines so wrap behavior is testable too.
    *[{"kind": "content",
       "out": f"test-content-0{i+1}-{mode_label}.pptx",
       "title": "The relationship model — without the relationship cost",
       "bullet_mode": mode,
       "items": [
           (0, "Three pillars work together to build the kind of wealth that compounds across decades"),
           (1, "Our integrated platform brings cash flow, balance sheet, and stewardship into one experience — eliminating the friction of moving between disconnected apps that never talk to each other."),
           (1, "Each pillar reinforces the others over time. Cash flow funds investments, investments grow your balance sheet, and stewardship protects what compounds."),
           (2, "Cash flow: earn more on your money while staying liquid"),
           (2, "Balance sheet: invest for the long term, borrow when it makes sense"),
           (2, "Stewardship: protect what you've built and optimize how it grows"),
           (3, "Dynamic Rebalancing keeps allocations on target"),
           (4, "Even during volatile market conditions"),
           (0, "Automation works on your terms, not against them — discipline becomes the platform's job"),
           (1, "Trade windows enforce discipline by design, protecting against emotional decisions during market volatility — exactly when you most need a steady hand on the wheel."),
           (2, "You set the strategy once"),
           (2, "M1 keeps you accountable to it"),
       ]}
      for i, (mode, mode_label) in enumerate([
          ("default",              "default"),
          ("text_with_bullets",    "text-with-bullets"),
          ("bullets_with_bullets", "bullets-with-bullets"),
      ])],
    # Big deck with many L1 groups — exercises:
    #   - Pagination across many pages
    #   - Image-slide cadence (1 in every 4 pages)
    #   - S10 (image left) cadence (every 3rd image slide)
    #   - Token-matched images from the OneDrive library
    # Each L1 has concrete, visualizable subheads so token matching has signal.
    {"kind": "content", "out": "test-content-05-imagery-cadence.pptx",
     "title": "M1 in everyday life",
     "bullet_mode": "default",
     "items": [
         (0, "Family budgeting on the kitchen computer"),
         (1, "Most households juggle competing priorities — groceries, bills, savings, the kids' activities — and the budget never quite balances on its own."),
         (2, "Track every account from one dashboard"),
         (2, "Forecast next month before it arrives"),
         (0, "Banking online while you celebrate the small wins"),
         (1, "M1 makes online banking feel like a tool that's actually on your side. No surprise fees, no friction, just visibility into where every dollar is going."),
         (2, "High-Yield Cash leads national averages"),
         (2, "Move money instantly between accounts"),
         (0, "Wealth strategy that respects your time"),
         (1, "Build a Pie that reflects your real allocation strategy, then let Dynamic Rebalancing keep it on target without forcing emotional trades."),
         (2, "Set the strategy once, refine quarterly"),
         (2, "Trade windows protect from market panic"),
         (0, "Baby's first steps toward generational wealth"),
         (1, "Custodial accounts let you start investing for the kids the day they're born, with the same Pies and Dynamic Rebalancing you use for yourself."),
         (2, "Custodial Pies grow with the child"),
         (2, "Tax-advantaged structures available"),
         (0, "401k rollover without the paperwork drama"),
         (1, "Bring your old workplace retirement account into M1 and unify your investing strategy across every dollar you have."),
         (2, "We handle the rollover paperwork"),
         (2, "Roll-in any 401k or IRA"),
         (0, "Reading the markets without the noise"),
         (1, "M1's dashboard shows you the signal that actually matters for long-term wealth-builders — without flashing tickers, hype, or doom-loop news."),
         (2, "Clean portfolio performance views"),
         (2, "Tax-loss harvest opportunities surfaced"),
         (0, "Beach chair vibes while the portfolio works"),
         (1, "Dynamic Rebalancing keeps your allocation on target whether you're at your desk or on vacation, so compounding never pauses."),
         (2, "Automation works on your timezone"),
         (2, "No need to log in during market hours"),
         (0, "Margin loans without selling the position"),
         (1, "Need liquidity? M1's Margin Loan lets you borrow against your portfolio at institutional rates without disrupting your long-term holdings."),
         (2, "Rates competitive with private banks"),
         (2, "Repay flexibly on your own schedule"),
     ]},
    # ── Comprehensive demo — every rule on one deck ──
    # Big enough to:
    #   • paginate across many slides (continuation titles)
    #   • trigger ≥3 image slides so the 3rd is S10 (image-left)
    #   • mix L1-L1 (+18pt) and L2-L2 (+12pt) group spacing
    #   • exercise the no-widow rule on the title
    {"kind": "content", "out": "demo-comprehensive.pptx",
     "title": "How M1 fits into real life",
     "bullet_mode": "default",
     "items": [
         (0, "A family kitchen budget that actually works"),
         (1, "Most households juggle bills, groceries, savings, and the kids' activities — and the budget never quite balances on its own. M1 brings them onto one screen."),
         (2, "Track every account from one dashboard"),
         (2, "Forecast next month before it arrives"),
         (3, "Visibility into where every dollar goes"),

         (0, "Online banking that celebrates the small wins"),
         (1, "High-Yield Cash beats brokerage sweep accounts by 3-5x for the same liquidity, and the gap compounds significantly over multi-year horizons."),
         (2, "Move money instantly between accounts"),
         (2, "Direct deposit ready in minutes"),

         (0, "Wealth strategy that respects your time"),
         (1, "Build a Pie that reflects your real allocation, then let Dynamic Rebalancing keep it on target without forced selling during volatility."),
         (1, "Trade windows enforce discipline by design, protecting against emotional decisions when markets test your conviction."),
         (2, "Set the strategy once, refine quarterly"),
         (2, "Automation runs on your terms"),

         (0, "Baby's first steps toward generational wealth"),
         (1, "Custodial Pies let you start investing for the kids the day they're born, with the same allocation framework you use yourself."),
         (2, "Custodial accounts grow with the child"),

         (0, "401k rollover without the paperwork drama"),
         (1, "Bring your old workplace retirement account into M1 and unify your investing strategy across every dollar you have."),
         (2, "We handle the rollover paperwork"),
         (2, "Roll-in any 401k or IRA"),

         (0, "Reading the markets without the noise"),
         (1, "M1's dashboard surfaces the signal that matters for long-term wealth-builders — no flashing tickers, no hype, no doom-loop news."),
         (2, "Clean portfolio performance views"),

         (0, "Beach chair vibes while the portfolio works"),
         (1, "Compounding does not require your attention. Dynamic Rebalancing keeps your allocation on target whether you're at your desk or on vacation."),
         (2, "Automation works on your timezone"),
         (2, "No need to log in during market hours"),
         (3, "Even during volatile market conditions"),

         (0, "Margin loans without selling the position"),
         (1, "Need liquidity? M1 lets you borrow against your portfolio at institutional rates without disrupting long-term holdings."),
         (2, "Rates competitive with private banks"),
         (2, "Repay flexibly on your own schedule"),

         (0, "Building a balance sheet that compounds"),
         (1, "M1 Pies express any allocation as a weighted tree, with fractional shares making every contribution efficient."),
         (2, "Pies grow alongside your goals"),

         (0, "Stewardship that protects what you've built"),
         (1, "Tax-loss harvesting integrated with portfolio strategy means optimization happens automatically — no quarterly scramble."),
         (2, "Aligns with your overall financial picture"),

         (0, "Dynamic Rebalancing without the manual labor"),
         (1, "When your portfolio drifts from target allocation, M1 redirects new contributions to underweight slices first — no forced selling, no taxable events."),
         (2, "Contributions correct drift automatically"),
         (2, "No reactive rebalancing during volatility"),

         (0, "Trade windows that protect long-term thinking"),
         (1, "Orders settle in scheduled windows once per market day, removing the temptation to react to intraday noise that erodes returns over time."),
         (2, "Discipline by design, not willpower"),

         (0, "Personal loans without the bank-branch friction"),
         (1, "M1 Personal Loans offer fixed rates from 7.99% APR with no origination fees, no late fees, and no prepayment penalties — fully digital."),
         (2, "Apply in minutes from your phone"),
         (2, "Funds available in days, not weeks"),

         (0, "Smart Transfers move money on your schedule"),
         (1, "Set rules that route your paycheck into savings, taxes, and investments automatically — every dollar gets a job before you see it."),
         (2, "Rules-based routing eliminates decision fatigue"),
         (2, "Adjusts as your income changes"),

         (0, "The dashboard that respects your attention"),
         (1, "M1's home screen surfaces what changed since you last logged in — no flashing tickers, no doom scroll, no emotional triggers."),
         (2, "Signal over noise, by design"),

         (0, "Retirement accounts that align with your strategy"),
         (1, "Traditional IRAs, Roth IRAs, and SEPs all use the same Pies framework as your taxable accounts — strategy stays consistent across every wrapper."),
         (2, "Tax-advantaged growth with familiar tooling"),
         (2, "Backdoor Roth and conversions supported"),

         (0, "Custodial accounts the kids never have to manage"),
         (1, "Open custodial brokerage accounts in their name with your same investment strategy — they take over at the age of majority with a built-in portfolio."),
         (2, "UGMA / UTMA accounts available"),

         (0, "Joint accounts for couples building together"),
         (1, "Combine accounts under a shared allocation strategy without losing individual visibility — perfect for households aligning on long-term goals."),
         (2, "Both partners see every position"),
         (2, "One Pie, two contributors"),

         (0, "Wealth that compounds across generations"),
         (1, "Estate planning integration helps you move assets efficiently to heirs while preserving your investment strategy across the transition."),

         (0, "Crypto access without the wallet headaches"),
         (1, "Buy and hold crypto alongside your equities in the same Pie framework — no separate exchange accounts, no key management to track."),
         (2, "Major coins available in custodial structure"),

         (0, "ESG-aligned investing without the lecture"),
         (1, "Express your values in your portfolio with vetted ESG Pies, or build a custom allocation that screens by your own criteria."),

         (0, "Tax documents that arrive when you need them"),
         (1, "Year-end 1099s and consolidated statements show up reliably each January with the cost-basis detail your accountant actually wants."),
         (2, "Direct download, no portal hunting"),

         (0, "Security that respects your peace of mind"),
         (1, "SIPC coverage on brokerage, FDIC on cash, two-factor on every login — the financial-grade security backbone you'd expect from a private bank."),
         (2, "Biometric login on mobile"),
         (2, "Activity alerts in real time"),

         (0, "Customer service that picks up the phone"),
         (1, "Reach a real human in minutes — no chatbot maze, no offshore call center. M1's support team is trained on the platform, not reading a script."),
         (2, "Available during market hours and beyond"),
     ]},
    # ── Multi-line title + imagery — exercises S11 routing ──
    # Long title forces 36pt S8 path. With enough content to paginate to 4+
    # pages, the image-slide is automatically routed to S11 (multi-line title +
    # image-right) instead of S9.
    {"kind": "content", "out": "demo-multiline-title-image.pptx",
     "title": "The relationship model — without the relationship cost: long-term wealth on your terms",
     "bullet_mode": "default",
     "items": [
         (0, "Three pillars work together to build the kind of wealth that compounds across decades"),
         (1, "Our integrated platform brings cash flow, balance sheet, and stewardship into one experience — eliminating friction between disconnected apps that never talk to each other."),
         (1, "Each pillar reinforces the others over time. Cash flow funds investments, investments grow your balance sheet, and stewardship protects what compounds."),
         (2, "Cash flow: earn more on your money while staying liquid"),
         (2, "Balance sheet: invest for the long term, borrow when it makes sense"),
         (2, "Stewardship: protect what you've built and optimize how it grows"),
         (3, "Dynamic Rebalancing keeps allocations on target"),
         (4, "Even during volatile market conditions"),

         (0, "Automation works on your terms, not against them"),
         (1, "Trade windows enforce discipline by design, protecting against emotional decisions during market volatility — exactly when you most need a steady hand."),
         (2, "You set the strategy once"),
         (2, "M1 keeps you accountable to it"),

         (0, "Margin loans give you liquidity without forced selling"),
         (1, "M1's Margin Loan lets you borrow against your portfolio at institutional rates, preserving the compounding engine that drives long-term wealth."),
         (2, "Rates competitive with institutional lenders"),
         (2, "No origination fees, no late fees"),

         (0, "High-Yield Cash earns more while you stay liquid"),
         (1, "Sweep idle cash into our high-yield account that consistently beats sweep accounts and money market funds — without sacrificing instant access."),
         (2, "Yields that lead national averages"),
     ]},
    # Fallback case — force the OneDrive catalog to look unavailable.
    # Expected: deck builds text-only (no Content + Image layouts) and one
    # warning is surfaced.
    {"kind": "content", "out": "test-content-06-no-onedrive-fallback.pptx",
     "title": "M1 in everyday life",
     "bullet_mode": "default",
     "env": {"M1_FORCE_NO_ONEDRIVE": "1"},
     "items": [
         (0, "Family budgeting on the kitchen computer"),
         (1, "Track every account from one dashboard, forecast next month before it arrives."),
         (2, "Visibility into where every dollar goes"),
         (2, "Forecast next month before it arrives"),
         (0, "Wealth strategy that respects your time"),
         (1, "Build a Pie that reflects your real allocation, then let Dynamic Rebalancing keep it on target."),
         (2, "Set the strategy once, refine quarterly"),
         (2, "Trade windows protect from market panic"),
         (0, "Banking online while you celebrate the small wins"),
         (1, "M1 makes online banking feel like a tool that's actually on your side."),
         (2, "High-Yield Cash leads national averages"),
         (2, "Move money instantly between accounts"),
         (0, "Margin loans without selling the position"),
         (1, "Need liquidity? Borrow against your portfolio at institutional rates without disrupting holdings."),
         (2, "Rates competitive with private banks"),
         (2, "Repay flexibly on your own schedule"),
     ]},
]


def build_case(case):
    out_path = os.path.join(OUT_DIR, case["out"])
    shutil.copy(TMPL, out_path)
    prs = Presentation(out_path)

    # Per-case env hook (e.g. force "no OneDrive" for the fallback test).
    # Clears module-level catalog cache so the next scan reflects the new env.
    case_env = case.get("env", {})
    prior_env = {}
    if case_env:
        pptx_helpers._onedrive_catalog_cache = None
        for k, v in case_env.items():
            prior_env[k] = os.environ.get(k)
            os.environ[k] = v

    if case["kind"] == "cover":
        slide = duplicate_slide(prs, COVER_SLIDE_INDEX)
        size, lines = set_cover_title(slide, case["title"])
    elif case["kind"] == "quote":
        slide = duplicate_slide(prs, QUOTE_SLIDE_INDEX)
        size, lines = set_quote(slide, case["quote"], case["attribution"])
    elif case["kind"] == "section":
        src_idx = SECTION_LIGHT_SLIDE_INDEX if case["mode"] == "light" else SECTION_DARK_SLIDE_INDEX
        slide = duplicate_slide(prs, src_idx)
        size, lines = set_section_break(slide, case["title"], case.get("supporting_text"))
    elif case["kind"] == "content":
        # Use the paginated builder — may produce multiple slides if the
        # content exceeds the body placeholder height.
        built = set_content_slides(prs, case["title"], case["items"],
                                   bullet_mode=case["bullet_mode"])
        slide = built[0]   # primary slide for verify_case width check
        # Stash diagnostic info for the report
        first_idx = pick_content_slide_index(case["title"])
        case["_picked_slide"] = "S7" if first_idx == 6 else "S8"
        case["_n_slides"] = len(built)
        size, lines = 0, [case["title"]]
    else:
        raise ValueError(f"Unknown case kind: {case['kind']}")

    for i in range(N_TEMPLATE_SLIDES - 1, -1, -1):
        delete_slide(prs, i)
    set_slide_view(prs)
    prs.save(out_path)

    # Stash any imagery warnings emitted during this build, then reset
    case["_imagery_warnings"] = list(pptx_helpers.imagery_warnings)

    # Restore env vars set above
    if case_env:
        for k, prev in prior_env.items():
            if prev is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = prev
        pptx_helpers._onedrive_catalog_cache = None

    return out_path, size, lines


def verify_case(case, out_path, size):
    prs = Presentation(out_path)
    slide = prs.slides[0]

    # Content slides: read body paragraphs at ph_idx=1, check level + bullet override
    if case["kind"] == "content":
        from lxml import etree
        nsa = "http://schemas.openxmlformats.org/drawingml/2006/main"
        body = next(p for p in slide.placeholders if p.placeholder_format.idx == 1)
        lines, line_info = [], []
        for para in body.text_frame.paragraphs:
            if not para.text.strip():
                continue
            pPr = para._element.find(f"{{{nsa}}}pPr")
            lvl = int(pPr.get('lvl', '0')) if pPr is not None else 0
            buChar = pPr.find(f"{{{nsa}}}buChar") if pPr is not None else None
            bu = buChar.get('char') if buChar is not None else "(inherited)"
            lines.append(para.text)
            line_info.append((lvl, bu, para.text))
        # Mode validation
        mode = case["bullet_mode"]
        expected_buChar = {
            "default":             {0: None, 1: None},
            "text_with_bullets":   {0: None, 1: "•"},
            "bullets_with_bullets":{0: "•",  1: "•"},
        }[mode]
        mode_ok = True
        for lvl, bu, _ in line_info:
            if lvl in expected_buChar:
                want = expected_buChar[lvl]
                if want is None and bu == "•":
                    mode_ok = False
                if want == "•" and bu != "•":
                    mode_ok = False
        content_info = " | ".join(f"L{lvl+1}({bu})" for lvl, bu, _ in line_info)
        return lines, [], 0, True, mode_ok, content_info

    ph = next(p for p in slide.placeholders if p.placeholder_format.idx == 0)
    ph_w_emu, _ = _get_ph_dims_from_xml(slide, 0)
    ph_w_pt = ph_w_emu / 12700
    lines = [p.text for p in ph.text_frame.paragraphs if p.text.strip()]

    line_widths = [(l, _text_width_pt(l, size)) for l in lines]
    fits = all(w <= ph_w_pt for _, w in line_widths)

    attr_ok = True
    if case["kind"] == "quote":
        attr_ph = next(p for p in slide.placeholders if p.placeholder_format.idx == 11)
        attr_ok = attr_ph.text_frame.text.startswith("— ")

    # For section breaks, also report supporting-text state
    section_info = None
    if case["kind"] == "section":
        support = _find_section_support_shape(slide)
        if support is None:
            section_info = "support shape: DELETED"
            section_info_ok = case.get("supporting_text") is None
        else:
            section_info = f"support shape: '{support.text_frame.text[:60]}'"
            section_info_ok = case.get("supporting_text") is not None
        attr_ok = section_info_ok

    return lines, line_widths, ph_w_pt, fits, attr_ok, section_info


def fmt_case(case):
    if case["kind"] == "cover":
        return f"cover    '{case['title']}'"
    if case["kind"] == "quote":
        return f"quote    '{case['quote']}' — {case['attribution']}"
    if case["kind"] == "section":
        st = case.get("supporting_text")
        st_label = f" + support" if st else " (no support)"
        return f"section/{case['mode']}{st_label}  '{case['title']}'"
    if case["kind"] == "content":
        picked = case.get("_picked_slide", "?")
        nsl = case.get("_n_slides", "?")
        return f"content/{case['bullet_mode']}  [{picked}, {nsl} slide(s)]  '{case['title']}'"
    return str(case)


def main():
    results = []
    for case in CASES:
        out_path, size, lines = build_case(case)
        actual, widths, ph_w_pt, fits, extra_ok, section_info = verify_case(case, out_path, size)

        print(f"\n── {case['out']} ─ {fmt_case(case)}")
        print(f"   {size}pt  {len(actual)} line(s)")
        for line, w in widths:
            status = "✓" if w <= ph_w_pt else "✗ OVERFLOW"
            print(f"     {w:6.1f} / {ph_w_pt:6.1f}pt  {status}  '{line}'")
        if section_info:
            print(f"   {section_info}")
        warns = case.get("_imagery_warnings", [])
        for w in warns:
            print(f"   ⚠ {w}")
        passed = fits and extra_ok
        print(f"   {'PASS' if passed else 'FAIL'}")
        results.append(passed)

    print()
    total, ok = len(results), sum(results)
    print(f"── Summary: {ok}/{total} cases passed ──")
    sys.exit(0 if ok == total else 1)


if __name__ == "__main__":
    main()
