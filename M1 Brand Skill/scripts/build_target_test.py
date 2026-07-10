#!/usr/bin/env python3
"""Targeted test for the S15 target graphic across 2, 3, and 4 levels.

Writes `testing/test-target-01.pptx` with one S15 slide per case:
  1. N=4 — the canonical case (Core → Need → Utility → Desire)
  2. N=3 — outermost ring + its pair textbox + label are removed
  3. N=2 — outermost two rings + pair textboxes + labels are removed

Verifies that:
  - Right-side subhead markers match the ring color of the corresponding
    concept (TOP marker = outermost / lightest, BOTTOM = innermost / darkest)
  - Unused outer rings/labels/pair-textboxes are removed (no placeholder
    'Subhead | Any key takeaways…' scaffolding remains)
  - Title routes through `set_content_title` (S8 multi-line treatment when
    needed)
"""
import os
import shutil
import sys

BASE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(BASE)
sys.path.insert(0, BASE)

from pptx import Presentation
from pptx_helpers import (
    duplicate_slide, delete_slide, routing_warnings, set_s15_target,
    set_slide_view, strip_helper_textboxes,
)
from render_for_review import render_slide_to_png

TMPL = os.path.join(ROOT, "assets", "M1-Presentation-Template_2026.pptx")
OUT  = os.path.join(ROOT, "testing", "test-target-01.pptx")

S15_TARGET = 14
N_TEMPLATE_SLIDES = 18

shutil.copy(TMPL, OUT)
prs = Presentation(OUT)

# Case 1 — N=4 (canonical). Verifies the "Foundation" label no longer wraps.
s = duplicate_slide(prs, S15_TARGET)
set_s15_target(s, "The wealth journey — four layers", [
    ("Foundation", "Emergency cash that earns more than it sits idle"),    # Core
    ("Growth",     "Invested portfolios that compound across decades"),    # Need
    ("Leverage",   "Borrowing power without forced selling"),              # Utility
    ("Legacy",     "Stewardship that protects what you've built"),          # Desire
])

# Case 2 — N=3 with default (enlarge=True). Rings scale up to fill the area.
s = duplicate_slide(prs, S15_TARGET)
set_s15_target(s, "How M1 builds three layers of capability", [
    ("Cash",     "High-yield account that earns more than the bank pays"),
    ("Invest",   "Portfolios that compound with intelligent rebalancing"),
    ("Borrow",   "Margin and personal loans at institutional rates"),
])

# Case 3 — N=2 with default (enlarge=True). Two rings fill the full footprint.
s = duplicate_slide(prs, S15_TARGET)
set_s15_target(s, "The two-layer trust model", [
    ("Custodial", "Your assets sit with the custodian — never on M1's books"),
    ("Insured",   "SIPC + FDIC coverage applied to the right account types"),
])

# Case 4 — N=2 with summary text. Rings stay at template positions; summary
# body text fills the freed area above.
s = duplicate_slide(prs, S15_TARGET)
set_s15_target(
    s,
    "The two-layer trust model with summary",
    [
        ("Custodial", "Your assets sit with the custodian — never on M1's books"),
        ("Insured",   "SIPC + FDIC coverage applied to the right account types"),
    ],
    summary="Your money never sits on M1's balance sheet — it stays with regulated custodians and is insured by SIPC and FDIC.",
)

# Delete originals, cleanup, save
for i in range(N_TEMPLATE_SLIDES - 1, -1, -1):
    delete_slide(prs, i)
strip_helper_textboxes(prs)
set_slide_view(prs)
prs.save(OUT)

print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
if routing_warnings:
    print(f"\nRouting warnings ({len(routing_warnings)}):")
    for w in routing_warnings:
        print(f"  ⚠ {w}")

# Render every slide for visual verification
print()
for i in range(len(prs.slides)):
    try:
        png = render_slide_to_png(OUT, i)
        print(f"Verify slide{i+1}: {png}")
    except Exception as e:
        print(f"⚠ Render failed for slide{i+1}: {e}")
