# -*- coding: utf-8 -*-
"""
Proven helper functions for M1 brand PowerPoint generation.
Tested 2026-04-16 — package-level cloning preserves all relationships (images, layouts).
"""

import copy
import math
import os
import re
import site
import sys

from PIL import Image as PILImage
from lxml import etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.parts.slide import SlidePart
from pptx.util import Emu

# ── Inter title-width metrics (loaded once, used for accurate line-width checks) ─
_usp = site.getusersitepackages()
if _usp not in sys.path:
    sys.path.insert(0, _usp)

_INTER_THIN_METRICS = None   # (cmap, hmtx, units_per_em) or None if unavailable

# Candidate font files for title width measurement, best first. The installed
# variable font (InterVariable.ttf) is the org-wide standard; the legacy static
# Inter-Thin.ttf is kept as a fallback for older machines. Measuring with the
# variable font's default (Regular) instance is conservative for thin titles
# (Regular is wider than Thin), so it never under-estimates width.
_INTER_METRIC_CANDIDATES = (
    "~/Library/Fonts/InterVariable.ttf",
    "~/Library/Fonts/Inter-Thin.ttf",
)

def _load_inter_thin():
    global _INTER_THIN_METRICS
    if _INTER_THIN_METRICS is not None:
        return _INTER_THIN_METRICS
    from fontTools.ttLib import TTFont
    for cand in _INTER_METRIC_CANDIDATES:
        try:
            font = TTFont(os.path.expanduser(cand))
            _INTER_THIN_METRICS = (
                font.getBestCmap(), font["hmtx"].metrics, font["head"].unitsPerEm,
            )
            return _INTER_THIN_METRICS
        except Exception:
            continue
    _INTER_THIN_METRICS = ()   # tried every candidate and failed
    return _INTER_THIN_METRICS

def _text_width_pt(text, size_pt):
    """Return the rendered advance-width of text in the brand Inter face at
    size_pt points. Falls back to a conservative heuristic (0.55× size per
    char) if no font file is available."""
    m = _load_inter_thin()
    if len(m) == 3:
        cmap, hmtx, upm = m
        total = sum(hmtx.get(cmap.get(ord(ch), ""), (int(upm * 0.55), 0))[0] for ch in text)
        return total / upm * size_pt
    return len(text) * size_pt * 0.55   # conservative fallback

nsmap = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

# ── M1 brand fonts — single source of truth ──────────────────
# Inter ships as a variable font (file: InterVariable.ttf, family name
# "Inter Variable"). PowerPoint resolves its weight named-instances to the
# family strings below — this is the org-wide standard, and it matches the
# Full Names the installed font exposes (e.g. "Inter Variable Thin"). Every
# title and body run, plus the template theme major/minor fonts, resolve to
# these names. To re-point the whole skill to a different cut, change these
# four values.
BRAND_FONT_TITLE = "Inter Variable Thin"         # headlines / slide titles
BRAND_FONT_BODY = "Inter Variable"               # body copy / default
BRAND_FONT_SEMIBOLD = "Inter Variable SemiBold"  # subheads / emphasis
BRAND_FONT_MEDIUM = "Inter Variable Medium"      # medium emphasis

# ── Safe placeholder text accessor ───────────────────────────

def get_ph_text(ph, fallback="<picture>"):
    """Safely get placeholder text. PlaceholderPicture has no .text attribute
    and will raise AttributeError — returns fallback string instead."""
    try:
        return ph.text
    except AttributeError:
        return fallback


# ── Smart quotes ──────────────────────────────────────────────
# Define curly quote chars as variables — do NOT use \u escapes in re.sub() replacements.
LDQ, RDQ, LSQ, RSQ = "\u201c", "\u201d", "\u2018", "\u2019"


def smart_quotes(text):
    """Replace straight quotes with curly quotes."""
    text = re.sub(r'"(\S)', LDQ + r'\1', text)
    text = re.sub(r'(\S)"', r'\1' + RDQ, text)
    text = text.replace('"', RDQ)
    text = re.sub(r"(\w)'(\w)", r'\1' + RSQ + r'\2', text)
    text = re.sub(r"'(\S)", LSQ + r'\1', text)
    text = re.sub(r"(\S)'", r'\1' + RSQ, text)
    text = text.replace("'", RSQ)
    return text


# ── Slide cloning (package-level) ────────────────────────────

def duplicate_slide(prs, slide_index):
    """Clone a slide preserving all relationships (images, layouts, etc)."""
    source_part = prs.slides[slide_index].part
    prs_part = prs.part

    # Find next available slide number
    existing_nums = []
    for part in prs_part.package.iter_parts():
        pn = str(part.partname)
        if pn.startswith("/ppt/slides/slide") and pn.endswith(".xml"):
            existing_nums.append(int(pn.replace("/ppt/slides/slide", "").replace(".xml", "")))
    next_num = max(existing_nums) + 1

    # Deep-copy XML and create new SlidePart
    new_xml = copy.deepcopy(source_part._element)
    new_slide_part = SlidePart(
        PackURI(f"/ppt/slides/slide{next_num}.xml"),
        source_part.content_type,
        source_part.package,
        new_xml,
    )

    # Copy all relationships (images, layout, etc) — skip notesSlide:
    # the notes slide contains a back-reference to its source slide, so copying
    # it causes the original slide to remain reachable after deletion (orphan in zip).
    _NOTES_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
    for rel in source_part.rels.values():
        if rel.reltype == _NOTES_RELTYPE:
            continue
        if rel.is_external:
            new_slide_part._rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_slide_part.relate_to(rel.target_part, rel.reltype)

    # Remap rIds in the XML (new rels may have different rIds)
    old_map = {r.rId: ("ext", r.reltype, r.target_ref) if r.is_external
               else ("int", r.reltype, r.target_part) for r in source_part.rels.values()}
    new_by_target = {}
    for r in new_slide_part.rels.values():
        key = ("ext", r.reltype, r.target_ref) if r.is_external else ("int", r.reltype, id(r.target_part))
        new_by_target[key] = r.rId
    rid_map = {}
    for old_rid, info in old_map.items():
        key = (info[0], info[1], info[2] if info[0] == "ext" else id(info[2]))
        if key in new_by_target:
            rid_map[old_rid] = new_by_target[key]

    r_ns = nsmap["r"]
    for elem in new_xml.iter():
        for attr in [f"{{{r_ns}}}embed", f"{{{r_ns}}}link", f"{{{r_ns}}}id"]:
            val = elem.get(attr)
            if val and val in rid_map:
                elem.set(attr, rid_map[val])

    # Register in presentation
    rId = prs_part.relate_to(new_slide_part, RT.SLIDE)
    sldIdLst = prs_part._element.find(f"{{{nsmap['p']}}}sldIdLst")
    max_id = max((int(s.get("id", "0")) for s in sldIdLst), default=255)
    new_sldId = etree.SubElement(sldIdLst, f"{{{nsmap['p']}}}sldId")
    new_sldId.set("id", str(max_id + 1))
    new_sldId.set(f"{{{nsmap['r']}}}id", rId)

    return new_slide_part.slide


# ── Create slide from any layout (including cross-master) ────

def find_layout_by_name(prs, name):
    """Find a slide layout part by name across all slide masters.
    python-pptx only exposes layouts from the first master via prs.slide_layouts.
    This searches all slideLayout parts in the package by parsing their XML.
    Returns the layout Part object, or raises ValueError.
    """
    for part in prs.part.package.iter_parts():
        pn = str(part.partname)
        if "slideLayout" in pn:
            xml = etree.fromstring(part.blob)
            cSld = xml.find(f"{{{nsmap['p']}}}cSld")
            if cSld is not None and cSld.get("name", "") == name:
                return part
    raise ValueError(f"Layout '{name}' not found in any slide master")


def add_slide_from_layout(prs, layout_name):
    """Create a new blank slide using any named layout (works across all masters).
    Use this instead of duplicate_slide() when the layout has no example slide
    in the deck (e.g. 'Content + Image' which exists only as a master layout).
    Returns the new Slide object.
    """
    layout_part = find_layout_by_name(prs, layout_name)
    prs_part = prs.part

    # Find next available slide number
    existing_nums = []
    for part in prs_part.package.iter_parts():
        pn = str(part.partname)
        if pn.startswith("/ppt/slides/slide") and pn.endswith(".xml"):
            existing_nums.append(int(pn.replace("/ppt/slides/slide", "").replace(".xml", "")))
    next_num = max(existing_nums) + 1

    # Build minimal slide XML using parse_xml so python-pptx registers custom element classes
    slide_xml = parse_xml(
        f'<p:sld xmlns:a="{nsmap["a"]}" xmlns:r="{nsmap["r"]}" xmlns:p="{nsmap["p"]}">'
        f'  <p:cSld><p:spTree>'
        f'    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        f'    <p:grpSpPr/>'
        f'  </p:spTree></p:cSld>'
        f'</p:sld>'
    )

    # Copy placeholder shapes from the layout into the slide
    layout_xml = etree.fromstring(layout_part.blob)
    layout_spTree = layout_xml.find(f"{{{nsmap['p']}}}cSld/{{{nsmap['p']}}}spTree")
    slide_spTree = slide_xml.find(f"{{{nsmap['p']}}}cSld/{{{nsmap['p']}}}spTree")
    for sp in layout_spTree.findall(f"{{{nsmap['p']}}}sp"):
        nvPr = sp.find(f"{{{nsmap['p']}}}nvSpPr/{{{nsmap['p']}}}nvPr/{{{nsmap['p']}}}ph")
        if nvPr is not None:
            slide_spTree.append(copy.deepcopy(sp))

    # Create SlidePart
    new_slide_part = SlidePart(
        PackURI(f"/ppt/slides/slide{next_num}.xml"),
        "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
        prs_part.package,
        slide_xml,
    )

    # Link to layout
    new_slide_part.relate_to(layout_part, RT.SLIDE_LAYOUT)

    # Copy layout's relationships that the slide might need (images in layout, etc.)
    for rel in layout_part.rels.values():
        if not rel.is_external and "slideMaster" not in rel.reltype:
            # Don't re-add layout rel types the slide doesn't need
            pass

    # Register in presentation
    rId = prs_part.relate_to(new_slide_part, RT.SLIDE)
    sldIdLst = prs_part._element.find(f"{{{nsmap['p']}}}sldIdLst")
    max_id = max((int(s.get("id", "0")) for s in sldIdLst), default=255)
    new_sldId = etree.SubElement(sldIdLst, f"{{{nsmap['p']}}}sldId")
    new_sldId.set("id", str(max_id + 1))
    new_sldId.set(f"{{{nsmap['r']}}}id", rId)

    return new_slide_part.slide


# ── Image insertion helpers ───────────────────────────────────

def _get_ph_dims_from_xml(slide, ph_idx):
    """Get placeholder dimensions from the slide XML or layout, in EMUs.
    Handles inherited dimensions when spPr/xfrm is absent on the slide shape.
    """
    # First try the slide's own shape
    for sp in slide._element.findall(f".//{{{nsmap['p']}}}sp"):
        ph_el = sp.find(f"{{{nsmap['p']}}}nvSpPr/{{{nsmap['p']}}}nvPr/{{{nsmap['p']}}}ph")
        if ph_el is not None and ph_el.get("idx", "0") == str(ph_idx):
            xfrm = sp.find(f"{{{nsmap['p']}}}spPr/{{{nsmap['a']}}}xfrm")
            if xfrm is not None:
                ext = xfrm.find(f"{{{nsmap['a']}}}ext")
                if ext is not None:
                    return int(ext.get("cx")), int(ext.get("cy"))

    # Fall back to layout
    layout_part = None
    for rel in slide.part.rels.values():
        if "slideLayout" in rel.reltype:
            layout_part = rel.target_part
            break
    if layout_part:
        layout_xml = etree.fromstring(layout_part.blob)
        for sp in layout_xml.findall(f".//{{{nsmap['p']}}}sp"):
            ph_el = sp.find(f"{{{nsmap['p']}}}nvSpPr/{{{nsmap['p']}}}nvPr/{{{nsmap['p']}}}ph")
            if ph_el is not None and ph_el.get("idx", "0") == str(ph_idx):
                xfrm = sp.find(f"{{{nsmap['p']}}}spPr/{{{nsmap['a']}}}xfrm")
                if xfrm is not None:
                    ext = xfrm.find(f"{{{nsmap['a']}}}ext")
                    if ext is not None:
                        return int(ext.get("cx")), int(ext.get("cy"))

    raise ValueError(f"Cannot determine dimensions for placeholder idx={ph_idx}")


def _get_ph_pos_from_xml(slide, ph_idx):
    """Get placeholder position (left, top) from slide XML or layout, in EMUs."""
    for sp in slide._element.findall(f".//{{{nsmap['p']}}}sp"):
        ph_el = sp.find(f"{{{nsmap['p']}}}nvSpPr/{{{nsmap['p']}}}nvPr/{{{nsmap['p']}}}ph")
        if ph_el is not None and ph_el.get("idx", "0") == str(ph_idx):
            xfrm = sp.find(f"{{{nsmap['p']}}}spPr/{{{nsmap['a']}}}xfrm")
            if xfrm is not None:
                off = xfrm.find(f"{{{nsmap['a']}}}off")
                if off is not None:
                    return int(off.get("x")), int(off.get("y"))

    layout_part = None
    for rel in slide.part.rels.values():
        if "slideLayout" in rel.reltype:
            layout_part = rel.target_part
            break
    if layout_part:
        layout_xml = etree.fromstring(layout_part.blob)
        for sp in layout_xml.findall(f".//{{{nsmap['p']}}}sp"):
            ph_el = sp.find(f"{{{nsmap['p']}}}nvSpPr/{{{nsmap['p']}}}nvPr/{{{nsmap['p']}}}ph")
            if ph_el is not None and ph_el.get("idx", "0") == str(ph_idx):
                xfrm = sp.find(f"{{{nsmap['p']}}}spPr/{{{nsmap['a']}}}xfrm")
                if xfrm is not None:
                    off = xfrm.find(f"{{{nsmap['a']}}}off")
                    if off is not None:
                        return int(off.get("x")), int(off.get("y"))

    raise ValueError(f"Cannot determine position for placeholder idx={ph_idx}")


def _set_ph_offset(slide, ph_idx, x_emu=None, y_emu=None):
    """Override a placeholder's slide-level position (x,y) in EMUs.

    When a placeholder has no slide-level `<a:xfrm>`, it inherits position
    and size from the layout. To shift the position on a single slide
    without affecting the layout, we have to write an `<a:xfrm>` with both
    `<a:off>` and `<a:ext>` (the size is needed so it doesn't fall back
    to layout's size with new position). Pass `x_emu=None` or `y_emu=None`
    to keep the layout-inherited coordinate for that axis."""
    for sp in slide._element.findall(f".//{{{nsmap['p']}}}sp"):
        ph_el = sp.find(f"{{{nsmap['p']}}}nvSpPr/{{{nsmap['p']}}}nvPr/{{{nsmap['p']}}}ph")
        if ph_el is None:
            continue
        if ph_el.get('idx', '0') != str(ph_idx):
            continue
        spPr = sp.find(f"{{{nsmap['p']}}}spPr")
        if spPr is None:
            spPr = etree.SubElement(sp, f"{{{nsmap['p']}}}spPr")
        xfrm = spPr.find(f"{{{nsmap['a']}}}xfrm")
        if xfrm is None:
            # Need to materialize xfrm with current layout-inherited dims
            ph_w, ph_h = _get_ph_dims_from_xml(slide, ph_idx)
            cur_x, cur_y = _get_ph_pos_from_xml(slide, ph_idx)
            xfrm = etree.SubElement(spPr, f"{{{nsmap['a']}}}xfrm")
            off = etree.SubElement(xfrm, f"{{{nsmap['a']}}}off")
            ext = etree.SubElement(xfrm, f"{{{nsmap['a']}}}ext")
            off.set('x', str(x_emu if x_emu is not None else cur_x))
            off.set('y', str(y_emu if y_emu is not None else cur_y))
            ext.set('cx', str(ph_w))
            ext.set('cy', str(ph_h))
        else:
            off = xfrm.find(f"{{{nsmap['a']}}}off")
            if off is None:
                cur_x, cur_y = _get_ph_pos_from_xml(slide, ph_idx)
                off = etree.SubElement(xfrm, f"{{{nsmap['a']}}}off")
                off.set('x', str(x_emu if x_emu is not None else cur_x))
                off.set('y', str(y_emu if y_emu is not None else cur_y))
            else:
                if x_emu is not None: off.set('x', str(x_emu))
                if y_emu is not None: off.set('y', str(y_emu))
        return True
    return False


def insert_picture(slide, ph_idx, image_path, mode="fit"):
    """Insert an image into a picture placeholder with fit or crop mode.

    mode="fit"  (default): Scale image to fit entirely within the placeholder.
                The full image is visible — no cropping. Centered within the
                placeholder bounds. Uses add_picture() with calculated size.
    mode="crop": Fill the placeholder completely, cropping the image as needed.
                Uses python-pptx's default insert_picture() behavior.
    """
    ph = None
    for p in slide.placeholders:
        if p.placeholder_format.idx == ph_idx:
            ph = p
            break
    if ph is None:
        raise ValueError(f"Placeholder idx={ph_idx} not found on slide")

    if mode == "crop":
        ph.insert_picture(image_path)
        return

    # ── Fit mode: add_picture() with calculated dimensions ──
    # Get placeholder bounds from XML (handles inherited dimensions)
    ph_w_emu, ph_h_emu = _get_ph_dims_from_xml(slide, ph_idx)
    ph_left, ph_top = _get_ph_pos_from_xml(slide, ph_idx)

    # Get image native dimensions
    with PILImage.open(image_path) as img:
        img_w, img_h = img.size
    img_aspect = img_w / img_h
    ph_aspect = ph_w_emu / ph_h_emu

    # Calculate fit dimensions (scale to constraining dimension)
    if img_aspect > ph_aspect:
        # Image is wider — constrain by width
        pic_w = ph_w_emu
        pic_h = int(ph_w_emu / img_aspect)
    else:
        # Image is taller — constrain by height
        pic_h = ph_h_emu
        pic_w = int(ph_h_emu * img_aspect)

    # Center within placeholder bounds
    left = ph_left + (ph_w_emu - pic_w) // 2
    top = ph_top + (ph_h_emu - pic_h) // 2

    # Remove the empty picture placeholder shape
    ph._element.getparent().remove(ph._element)

    # Add the picture at calculated position and size
    slide.shapes.add_picture(image_path, left, top, pic_w, pic_h)


# ── Slide deletion ────────────────────────────────────────────

def delete_slide(prs, slide_index):
    """Delete a slide by index.
    Drops the relationship AND removes the part from the package cache so no
    orphaned slide*.xml files remain in the saved zip (which causes PowerPoint
    to show a 'found a problem' repair prompt on open).
    """
    slide_part = prs.slides[slide_index].part
    prs_part = prs.part
    rId = next(r.rId for r in prs_part.rels.values()
               if not r.is_external and r.target_part is slide_part)
    sldIdLst = prs_part._element.find(f"{{{nsmap['p']}}}sldIdLst")
    for sldId in list(sldIdLst):
        if sldId.get(f"{{{nsmap['r']}}}id") == rId:
            sldIdLst.remove(sldId)
            break
    prs_part.drop_rel(rId)


# ── Text helpers ──────────────────────────────────────────────

def set_ph(slide, ph_idx, text):
    """Set single-line placeholder text, preserving run formatting."""
    text = smart_quotes(text)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = text
                for run in tf.paragraphs[0].runs[1:]:
                    run.text = ""
            else:
                tf.paragraphs[0].text = text
            while len(tf.paragraphs) > 1:
                tf.paragraphs[-1]._element.getparent().remove(tf.paragraphs[-1]._element)
            return True
    return False


def _make_run(parent, text, size_pt, font_name=BRAND_FONT_BODY):
    """Create an <a:r> run element with font and size."""
    text = smart_quotes(text)
    r = etree.SubElement(parent, f"{{{nsmap['a']}}}r")
    rPr = etree.SubElement(r, f"{{{nsmap['a']}}}rPr")
    rPr.set("lang", "en-US")
    rPr.set("sz", str(size_pt * 100))
    rPr.set("dirty", "0")
    latin = etree.SubElement(rPr, f"{{{nsmap['a']}}}latin")
    latin.set("typeface", font_name)
    t = etree.SubElement(r, f"{{{nsmap['a']}}}t")
    t.text = text


# ── Agenda (numbered list for empty placeholders) ────────────

def fill_agenda(placeholder, items, font_name=BRAND_FONT_BODY, font_size_pt=18):
    """Build numbered agenda with optional bullet subsections.
    items: list of strings or dicts {"text": "...", "subs": ["..."]}.
    """
    txBody = placeholder._element.find(f"{{{nsmap['p']}}}txBody")
    if txBody is None:
        txBody = placeholder._element.find(f"{{{nsmap['a']}}}txBody")
    for p in txBody.findall(f"{{{nsmap['a']}}}p"):
        txBody.remove(p)

    for item in items:
        text = item if isinstance(item, str) else item["text"]
        subs = [] if isinstance(item, str) else item.get("subs", [])

        # Numbered top-level item
        p = etree.SubElement(txBody, f"{{{nsmap['a']}}}p")
        pPr = etree.SubElement(p, f"{{{nsmap['a']}}}pPr")
        pPr.set("lvl", "0")
        buAutoNum = etree.SubElement(pPr, f"{{{nsmap['a']}}}buAutoNum")
        buAutoNum.set("type", "arabicPeriod")
        spcBef = etree.SubElement(pPr, f"{{{nsmap['a']}}}spcBef")
        spcPts = etree.SubElement(spcBef, f"{{{nsmap['a']}}}spcPts")
        spcPts.set("val", "600")
        _make_run(p, text, font_size_pt, font_name)

        # Optional bullet subsections (indented)
        for sub in subs:
            sp = etree.SubElement(txBody, f"{{{nsmap['a']}}}p")
            spPr = etree.SubElement(sp, f"{{{nsmap['a']}}}pPr")
            spPr.set("lvl", "1")
            buChar = etree.SubElement(spPr, f"{{{nsmap['a']}}}buChar")
            buChar.set("char", "\u2022")
            _make_run(sp, sub, font_size_pt - 2, font_name)


# ── Bullet list (for empty placeholders) ─────────────────────

def fill_bullets(placeholder, items, font_name=BRAND_FONT_BODY, font_size_pt=18):
    """Build bullet paragraphs from scratch for empty placeholders."""
    txBody = placeholder._element.find(f"{{{nsmap['p']}}}txBody")
    if txBody is None:
        txBody = placeholder._element.find(f"{{{nsmap['a']}}}txBody")
    for p in txBody.findall(f"{{{nsmap['a']}}}p"):
        txBody.remove(p)
    for item in items:
        p = etree.SubElement(txBody, f"{{{nsmap['a']}}}p")
        pPr = etree.SubElement(p, f"{{{nsmap['a']}}}pPr")
        pPr.set("lvl", "0")
        buChar = etree.SubElement(pPr, f"{{{nsmap['a']}}}buChar")
        buChar.set("char", "\u2022")
        r = etree.SubElement(p, f"{{{nsmap['a']}}}r")
        rPr = etree.SubElement(r, f"{{{nsmap['a']}}}rPr")
        rPr.set("lang", "en-US")
        rPr.set("sz", str(font_size_pt * 100))
        rPr.set("dirty", "0")
        latin = etree.SubElement(rPr, f"{{{nsmap['a']}}}latin")
        latin.set("typeface", font_name)
        t = etree.SubElement(r, f"{{{nsmap['a']}}}t")
        t.text = smart_quotes(item)


# ── Bullet list (for placeholders with existing runs) ────────

def set_ph_bullets(slide, ph_idx, items):
    """Set bullets in a placeholder that has existing runs to clone."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            template_p = copy.deepcopy(tf.paragraphs[0]._element)
            txBody = tf._txBody
            for p in list(txBody.findall(f"{{{nsmap['a']}}}p")):
                txBody.remove(p)
            for item in items:
                item = smart_quotes(item)
                new_p = copy.deepcopy(template_p)
                runs = new_p.findall(f"{{{nsmap['a']}}}r")
                if runs:
                    runs[0].find(f"{{{nsmap['a']}}}t").text = item
                    for r in runs[1:]:
                        r.find(f"{{{nsmap['a']}}}t").text = ""
                txBody.append(new_p)
            return True
    return False


# ── Label/description split bullets ──────────────────────────

def set_ph_split_bullets(slide, ph_idx, items, font_name=BRAND_FONT_BODY):
    """Set bullets with label/description splitting.
    items: list of (label, description) tuples or plain strings.
    Tuples produce a top-level label + indented sub-bullet description.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame
            template_p = copy.deepcopy(tf.paragraphs[0]._element)
            txBody = tf._txBody
            for p in list(txBody.findall(f"{{{nsmap['a']}}}p")):
                txBody.remove(p)
            for item in items:
                if isinstance(item, tuple):
                    label, desc = item
                    # Top-level: label
                    np = copy.deepcopy(template_p)
                    runs = np.findall(f"{{{nsmap['a']}}}r")
                    if runs:
                        runs[0].find(f"{{{nsmap['a']}}}t").text = smart_quotes(label)
                        for r in runs[1:]:
                            r.find(f"{{{nsmap['a']}}}t").text = ""
                    txBody.append(np)
                    # Sub-bullet: description
                    sp = etree.SubElement(txBody, f"{{{nsmap['a']}}}p")
                    spPr = etree.SubElement(sp, f"{{{nsmap['a']}}}pPr")
                    spPr.set("lvl", "1")
                    buChar = etree.SubElement(spPr, f"{{{nsmap['a']}}}buChar")
                    buChar.set("char", "\u2022")
                    _make_run(sp, desc, 16, font_name)
                else:
                    np = copy.deepcopy(template_p)
                    runs = np.findall(f"{{{nsmap['a']}}}r")
                    if runs:
                        runs[0].find(f"{{{nsmap['a']}}}t").text = smart_quotes(item)
                        for r in runs[1:]:
                            r.find(f"{{{nsmap['a']}}}t").text = ""
                    txBody.append(np)
            return True
    return False


# ── Grouped bullet list (parent + indented children) ────────

def set_ph_grouped_bullets(slide, ph_idx, groups, group_space_pt=12):
    """Set grouped bullet content with master-slide styling and group spacing.

    groups: list of items. Each item is either:
      - str: plain parent-level item (no children)
      - tuple (str, list[str]): (parent_text, [child_text, ...])

    Approach — two passes:
      Pass 1 (style): Clone template paragraphs at the correct levels so that
        master slide formatting (bullet char, indent, color, size) is preserved
        exactly — no overrides:
          Parents  → clone lvl=0 template paragraph (navy bold, no bullet)
          Children → clone lvl=2 template paragraph (teal 14pt, master bullet + indent)
      Pass 2 (spacing): Walk the written paragraphs and apply spcBef on each
        group's first (parent) paragraph after the first group.

    NOTE on levels — the Content Text layout defines:
      lvl=0: subhead style (navy bold, buNone)
      lvl=1: secondary text (teal 16pt, buNone — NO bullet, use only for non-bulleted lines)
      lvl=2: bullet level (teal 14pt, inherits bullet char + indent from master)
    Never use lvl=1 for bulleted children — it explicitly suppresses the bullet.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            tf = ph.text_frame

            # ── Collect template paragraphs at lvl=0 and lvl=2 ──
            tmpl_p0 = None
            tmpl_p2 = None
            for para in tf.paragraphs:
                pPr = para._element.find(f"{{{nsmap['a']}}}pPr")
                lvl = int(pPr.get("lvl", "0")) if pPr is not None else 0
                if lvl == 0 and tmpl_p0 is None:
                    tmpl_p0 = copy.deepcopy(para._element)
                elif lvl == 2 and tmpl_p2 is None:
                    tmpl_p2 = copy.deepcopy(para._element)

            txBody = tf._txBody
            for p in list(txBody.findall(f"{{{nsmap['a']}}}p")):
                txBody.remove(p)

            # ── Pass 1: place text using cloned template paragraphs ──
            group_starts = []  # indices of each group's parent paragraph
            para_count = 0

            for item in groups:
                parent_text = item if isinstance(item, str) else item[0]
                children = [] if isinstance(item, str) else (item[1] if len(item) > 1 else [])

                # Parent — clone lvl=0
                parent_p = copy.deepcopy(tmpl_p0)
                runs = parent_p.findall(f"{{{nsmap['a']}}}r")
                if runs:
                    runs[0].find(f"{{{nsmap['a']}}}t").text = smart_quotes(parent_text)
                    for r in runs[1:]:
                        r.find(f"{{{nsmap['a']}}}t").text = ""
                txBody.append(parent_p)
                group_starts.append(para_count)
                para_count += 1

                # Children — clone lvl=2 (master bullet + indent)
                for child_text in children:
                    if tmpl_p2 is not None:
                        child_p = copy.deepcopy(tmpl_p2)
                        runs = child_p.findall(f"{{{nsmap['a']}}}r")
                        if runs:
                            runs[0].find(f"{{{nsmap['a']}}}t").text = smart_quotes(child_text)
                            for r in runs[1:]:
                                r.find(f"{{{nsmap['a']}}}t").text = ""
                        else:
                            _make_run(child_p, child_text, 14)
                    else:
                        # Fallback: build lvl=2 paragraph if template didn't have one
                        child_p = etree.SubElement(txBody, f"{{{nsmap['a']}}}p")
                        cpPr = etree.SubElement(child_p, f"{{{nsmap['a']}}}pPr")
                        cpPr.set("lvl", "2")
                        _make_run(child_p, child_text, 14)
                    txBody.append(child_p)
                    para_count += 1

            # ── Pass 2: apply group spacing on every parent after the first ──
            written = txBody.findall(f"{{{nsmap['a']}}}p")
            for idx in group_starts[1:]:
                p_el = written[idx]
                pPr = p_el.find(f"{{{nsmap['a']}}}pPr")
                if pPr is None:
                    pPr = etree.Element(f"{{{nsmap['a']}}}pPr")
                    p_el.insert(0, pPr)
                spcBef = pPr.find(f"{{{nsmap['a']}}}spcBef")
                if spcBef is None:
                    spcBef = etree.SubElement(pPr, f"{{{nsmap['a']}}}spcBef")
                for child in list(spcBef):
                    spcBef.remove(child)
                spcPts = etree.SubElement(spcBef, f"{{{nsmap['a']}}}spcPts")
                spcPts.set("val", str(group_space_pt * 100))

            return True
    return False


# ── 5-level content bullets (S7 / S8) ─────────────────────────

CONTENT_BULLET_MODES = ("default", "text_with_bullets", "bullets_with_bullets")

# Per-level font sizes and left-margins (pt), inferred from the layout's
# lstStyle override on top of the master bodyStyle.
LEVEL_FONT_PT = {0: 18, 1: 16, 2: 14, 3: 12, 4: 11}
LEVEL_MARL_PT = {0: 0,  1: 0,  2: 45, 3: 63, 4: 81}

# Forced-bullet marL (pt) for L1/L2 in modes 2/3 — matches what `_force_bullet`
# stamps onto the paragraph's pPr.
L1_FORCED_BULLET_MARL_PT = 18
L2_FORCED_BULLET_MARL_PT = 27

# Inter Variable Thin metrics under-measure body text (which renders in a slightly wider
# weight); pad estimated widths by this factor so pagination errs on the safe
# side and we don't push content into the last 0.2" before the footer.
CONTENT_WIDTH_SAFETY_FACTOR = 1.15

# Line-height multiplier used specifically for content-slide body text. Larger
# than the cover-title `LINE_HEIGHT_FACTOR` (1.2) because: (a) body uses Inter
# Regular which has looser metrics than Inter Variable Thin, (b) PowerPoint adds implicit
# inter-paragraph spacing that isn't visible in the slide XML, and (c) we'd
# rather paginate one paragraph early than allow any footer overflow.
LINE_HEIGHT_FACTOR_CONTENT = 1.5

# Use only this fraction of the body placeholder's stated height as the
# pagination budget — extra safety buffer so pagination triggers BEFORE the
# 0.2" footer breathing room is touched.
CONTENT_USABLE_HEIGHT_FRACTION = 0.92

# Group-separator spacing (extra space-before applied between consecutive
# paragraphs at the same structural level — visually groups things that go
# together). Hoisted to module scope so pagination height estimates can apply
# them too.
L1_GROUP_SPCBEF_PT = 18   # between two L1 (subhead) paragraphs
L2_GROUP_SPCBEF_PT = 12   # between two L2 (body) paragraphs

# Content slide templates: S7 is for single-line titles (42pt default), S8 is
# for multi-line titles (36pt slide-level override, body shifted down). Pick by
# measuring whether the title fits in one line at 42pt.
CONTENT_S7_SLIDE_INDEX = 6        # S7 — single-line title
CONTENT_S8_SLIDE_INDEX = 7        # S8 — multi-line title
CONTENT_TITLE_WIDTH_PT = 864.0    # both S7 and S8 title placeholder width
CONTENT_S7_TITLE_PT    = 42       # S7 title default font size


def pick_content_slide_index(title):
    """Return the slide index to clone for an S7-style content slide:
      • S7 (index 6) when the title fits on one line at 42pt
      • S8 (index 7) when the title needs to wrap to multiple lines
    Uses the standard `SAFETY_PT` width margin so we don't sit on the edge."""
    fits_one_line = _text_width_pt(title, CONTENT_S7_TITLE_PT) <= CONTENT_TITLE_WIDTH_PT - SAFETY_PT
    return CONTENT_S7_SLIDE_INDEX if fits_one_line else CONTENT_S8_SLIDE_INDEX


CONTENT_S8_TITLE_MAX_PT = 36       # S8 default title size (slide-level rPr)
CONTENT_S8_TITLE_MIN_PT = 28       # floor for multi-line title size-down

# Exact S8 multi-line spec — extracted from the template:
# Title placeholder shifts from y=0 (S7) to y=246165 EMU (0.269") on S8.
# This is the "exact solution" that must be applied to every slide that
# inherits an S7-style title placeholder (S9, S10, S12, S13, S14, S15, S16,
# S17) when the title wraps to 2 lines — the y-shift gives 2 lines of 36pt
# Inter Variable Thin breathing room above the body anchor point.
S8_TITLE_Y_OFFSET_EMU = 246165


def set_content_title(slide, title, ph_idx=0):
    """Set the title on a content / image / chart / graphic slide. ALWAYS uses
    Inter Variable Thin (brand rule). When the title fits one line at 42pt it's left
    at the slide's default position. When the title needs to wrap, the EXACT
    S8 spec is applied: title placeholder shifted to y=0.269" (matches S8) AND
    title rendered at 36pt with phrase-aware breaks + no widows.

    The y-shift is critical — without it, 2-line titles on S7-cloned slides
    (S9, S12-S17) crowd the very top of the slide because their title
    placeholder starts at y=0 with bottom-anchored text.

    Returns (size_pt, lines).
    """
    ph_w_emu, _ = _get_ph_dims_from_xml(slide, ph_idx)
    ph_w_pt = ph_w_emu / 12700

    if _text_width_pt(title, CONTENT_S7_TITLE_PT) <= ph_w_pt - SAFETY_PT:
        # Single-line at 42pt (S7 path) — no position shift
        _write_styled_lines(slide, ph_idx, [title], CONTENT_S7_TITLE_PT)
        return CONTENT_S7_TITLE_PT, [title]

    # Multi-line — exact S8 treatment:
    #   1. Shift the title placeholder y down to match S8's 0.269" offset
    #   2. Render at 36pt (sizes down to 28pt if needed for very long titles)
    #   3. Phrase-aware breaks + widow rejection (already inside _layout_phrase_lines)
    _set_ph_offset(slide, ph_idx, y_emu=S8_TITLE_Y_OFFSET_EMU)
    size_pt, lines = _layout_phrase_lines(
        title, ph_w_pt,
        min_pt=CONTENT_S8_TITLE_MIN_PT, max_pt=CONTENT_S8_TITLE_MAX_PT,
    )
    _write_styled_lines(slide, ph_idx, lines, size_pt)
    return size_pt, lines


def set_content_bullets(slide, ph_idx, items, bullet_mode="default",
                        first_l1_continues=False):
    """Populate an S7 content placeholder with hierarchical content across the
    template's 5 styled levels (matches the screenshot reference):

        L1 (lvl=0): subhead          — navy bold, no bullet
        L2 (lvl=1): body text        — teal, no bullet
        L3 (lvl=2): bullet           — teal, • bullet
        L4 (lvl=3): sub-bullet       — light blue, - dash
        L5 (lvl=4): final sub-bullet — light blue, › chevron

    `items` is a list of (level, text) tuples where level is 0–4.

    `bullet_mode` controls whether L1 and L2 get bullets added on top of their
    default text styling:
        "default"              — L1 plain, L2 plain (recommended)
        "text_with_bullets"    — L1 plain, L2 gets • bullet
        "bullets_with_bullets" — L1 and L2 both get • bullets (not recommended)

    L3–L5 always carry their template bullets regardless of mode.

    `first_l1_continues=True` — used on continuation slides built by
    `set_content_slides`. Makes the very first L1 paragraph on the slide get
    the +18pt L1-group space-before (as if a previous L1 had been on the
    same page), so the L1↔L1 rhythm is preserved across the page break.

    Returns True on success, False if the placeholder isn't found.
    """
    if bullet_mode not in CONTENT_BULLET_MODES:
        raise ValueError(
            f"bullet_mode must be one of {CONTENT_BULLET_MODES}, got {bullet_mode!r}"
        )

    ph = None
    for p in slide.placeholders:
        if p.placeholder_format.idx == ph_idx:
            ph = p
            break
    if ph is None:
        return False

    tf = ph.text_frame
    txBody = tf._txBody

    # Collect one template paragraph per level from the existing slide content
    templates = {}
    for para in tf.paragraphs:
        pPr_el = para._element.find(f"{{{nsmap['a']}}}pPr")
        lvl = int(pPr_el.get('lvl', '0')) if pPr_el is not None else 0
        if lvl not in templates:
            templates[lvl] = copy.deepcopy(para._element)

    # Clear existing paragraphs
    for p in list(txBody.findall(f"{{{nsmap['a']}}}p")):
        txBody.remove(p)

    # Hanging-indent geometry for forced L1 / L2 bullets (EMU; 914400 = 1 inch).
    # The layout's lstStyle sets marL=0, indent=0 on lvl1/lvl2pPr (no bullet by
    # default), so when we force a bullet we must also set marL+indent or the
    # bullet character collides with the text. Indent here is more generous than
    # the master's -114277 (0.125") so the gap between • and text is visible.
    L1_MARL,   L1_INDENT   = 228600,  -285750    # L1 with forced bullet
    L2_MARL,   L2_INDENT   = 342900,  -285750    # L2 with forced bullet

    def _set_space_before(p_el, pt):
        """Set <a:spcBef><a:spcPts val="..."/></a:spcBef> on the paragraph's
        pPr, overriding any existing spcBef. Units: hundredths of a point
        (val="1800" = 18pt)."""
        pPr = p_el.find(f"{{{nsmap['a']}}}pPr")
        if pPr is None:
            pPr = etree.Element(f"{{{nsmap['a']}}}pPr")
            p_el.insert(0, pPr)
        for existing in pPr.findall(f"{{{nsmap['a']}}}spcBef"):
            pPr.remove(existing)
        spcBef = etree.SubElement(pPr, f"{{{nsmap['a']}}}spcBef")
        spcPts = etree.SubElement(spcBef, f"{{{nsmap['a']}}}spcPts")
        spcPts.set("val", str(int(pt * 100)))

    def _force_bullet(p_el, level, char="•"):
        """Override the paragraph's bullet to a • character AND set marL/indent
        so the bullet has proper hanging-indent spacing on a level whose layout
        normally has no bullet."""
        pPr = p_el.find(f"{{{nsmap['a']}}}pPr")
        if pPr is None:
            pPr = etree.Element(f"{{{nsmap['a']}}}pPr")
            p_el.insert(0, pPr)
        for tag in ("buNone", "buChar", "buAutoNum"):
            for existing in pPr.findall(f"{{{nsmap['a']}}}{tag}"):
                pPr.remove(existing)
        if level == 0:
            pPr.set("marL",   str(L1_MARL))
            pPr.set("indent", str(L1_INDENT))
        else:  # level == 1
            pPr.set("marL",   str(L2_MARL))
            pPr.set("indent", str(L2_INDENT))
        buChar = etree.SubElement(pPr, f"{{{nsmap['a']}}}buChar")
        buChar.set("char", char)

    # L1 spacing rule: ANY L1 after the first L1 ever (across slides) gets
    # +18pt. Track this with `saw_l1`, separate from `prev_level` (which is
    # used only for the L2↔L2 consecutive-paragraph check).
    saw_l1 = first_l1_continues
    prev_level = None
    for level, text in items:
        if level not in templates:
            raise ValueError(
                f"No template paragraph at level {level} in placeholder idx={ph_idx}"
            )
        new_p = copy.deepcopy(templates[level])

        # Replace text in the first run; clear the rest
        runs = new_p.findall(f"{{{nsmap['a']}}}r")
        if runs:
            runs[0].find(f"{{{nsmap['a']}}}t").text = smart_quotes(text)
            for r in runs[1:]:
                r.find(f"{{{nsmap['a']}}}t").text = ""
        else:
            _make_run(new_p, text, 16)

        # Mode-driven bullet override on L1 / L2
        if level == 1 and bullet_mode in ("text_with_bullets", "bullets_with_bullets"):
            _force_bullet(new_p, level=1)
        elif level == 0 and bullet_mode == "bullets_with_bullets":
            _force_bullet(new_p, level=0)

        # Group-separator spacing
        # - L1 gets +18pt if any L1 has appeared before it (including
        #   continuation from a prior slide)
        # - L2 gets +12pt only if the immediately previous paragraph on this
        #   slide was also L2 (consecutive body lines)
        if level == 0 and saw_l1:
            _set_space_before(new_p, L1_GROUP_SPCBEF_PT)
        elif level == 1 and prev_level == 1:
            _set_space_before(new_p, L2_GROUP_SPCBEF_PT)

        txBody.append(new_p)
        if level == 0:
            saw_l1 = True
        prev_level = level

    return True


# ── Content pagination (no overflow past footer, never split L1) ──────

# Body of the content placeholder must end at least 0.2" before the footer.
# Body placeholder bottom on both S7 and S8 already sits at ~6.75" while the
# footer line / slide-number is at ~6.95" — so the placeholder height IS the
# usable height. We rely on `_get_ph_dims_from_xml(slide, 1)` to read it.

def _item_height_pt(level, text, ph_w_pt, prev_level, bullet_mode, saw_l1=False):
    """Estimate the rendered height (pt) of a single content item paragraph,
    applying the same L1/L2 spacing rules as `set_content_bullets`:
      - L1 gets +18pt spcBef if `saw_l1` is True (any L1 has come before it)
      - L2 gets +12pt spcBef if `prev_level == 1` (consecutive L2 on same slide)
    """
    font_pt = LEVEL_FONT_PT[level]
    line_h  = font_pt * LINE_HEIGHT_FACTOR_CONTENT

    marL_pt = LEVEL_MARL_PT[level]
    if level == 0 and bullet_mode == "bullets_with_bullets":
        marL_pt = L1_FORCED_BULLET_MARL_PT
    elif level == 1 and bullet_mode in ("text_with_bullets", "bullets_with_bullets"):
        marL_pt = L2_FORCED_BULLET_MARL_PT
    eff_w_pt = max(1.0, ph_w_pt - marL_pt)

    text_w = _text_width_pt(text, font_pt) * CONTENT_WIDTH_SAFETY_FACTOR
    n_lines = max(1, math.ceil(text_w / eff_w_pt))

    spcBef = 0
    if level == 0 and saw_l1:
        spcBef = L1_GROUP_SPCBEF_PT
    elif level == 1 and prev_level == 1:
        spcBef = L2_GROUP_SPCBEF_PT

    return spcBef + n_lines * line_h


def _paginate_content_items(items, ph_w_pt, ph_h_pt, bullet_mode):
    """Split `items` into pages so the estimated total height of each page is
    ≤ `ph_h_pt`. Never breaks an L1 group across pages unless a single group
    is itself taller than one page (last-resort split — group goes on its own
    page and may visually overflow).

    Returns a list of `(page_items, first_l1_continues)` tuples — one per
    page — where `first_l1_continues` tells the body builder whether to apply
    the L1↔L1 spacing to the first L1 on that page.
    """
    # Group by L1 boundary
    groups, cur = [], []
    for item in items:
        if item[0] == 0 and cur:
            groups.append(cur)
            cur = []
        cur.append(item)
    if cur:
        groups.append(cur)

    def _group_height(group, is_first_group_overall):
        """Estimate height of one L1 group with proper L1/L2 spacing.
        `is_first_group_overall=True` means this is the very first group of
        the outline, so its leading L1 gets no spcBef. Subsequent groups'
        leading L1 always gets +18pt (matches the saw_l1 semantics)."""
        h = 0
        prev = None
        saw_l1 = not is_first_group_overall   # subsequent groups: pretend we've seen an L1
        for lvl, txt in group:
            h += _item_height_pt(lvl, txt, ph_w_pt, prev, bullet_mode, saw_l1)
            if lvl == 0:
                saw_l1 = True
            prev = lvl
        return h

    pages = []                     # list of (items_list, continues_flag)
    cur_items, cur_h = [], 0.0

    for gi, group in enumerate(groups):
        is_first_overall = (gi == 0)
        gh = _group_height(group, is_first_overall)

        if cur_items and cur_h + gh > ph_h_pt:
            pages.append((cur_items, len(pages) > 0))
            cur_items, cur_h = [], 0.0

        cur_items.extend(group)
        cur_h += gh

    if cur_items:
        pages.append((cur_items, len(pages) > 0))

    # Mark continuation flag: the very first page is never a continuation;
    # every page after the first IS a continuation (relative to the original
    # outline's first L1).
    pages = [(items_, idx > 0) for idx, (items_, _) in enumerate(pages)]
    return pages


def set_content_slides(prs, title, items, bullet_mode="default"):
    """Build one or more content slides for `title`/`items`/`bullet_mode`,
    paginating to keep content within the body placeholder height (which is
    sized to leave 0.2" of breathing room above the footer). Never splits an
    L1 group across slides unless a single group exceeds one page on its own.

    Continuation slides get ` (continued)` appended to the title. The first
    L1 on each continuation slide receives the L1↔L1 +18pt space-before so
    the visual rhythm of the original outline is preserved.

    Returns the list of slides built (1 if everything fits, 2+ if paginated).
    """
    # Reset per-build imagery warnings
    imagery_warnings.clear()

    # ── Pagination math (read S7's body dims without cloning) ──
    s7_tmpl = prs.slides[CONTENT_S7_SLIDE_INDEX]
    body_w_emu, body_h_emu = _get_ph_dims_from_xml(s7_tmpl, 1)
    body_w_pt = body_w_emu / 12700
    body_h_pt = body_h_emu / 12700
    budget_h_pt = body_h_pt * CONTENT_USABLE_HEIGHT_FRACTION

    # ── Pagination pass: text-only pages first ──
    pages = _paginate_content_items(items, body_w_pt, budget_h_pt, bullet_mode)

    # ── Imagery pass: decide which pages get images ──
    catalog = _scan_onedrive_image_catalog()
    image_for_page = {}                 # page_idx → absolute image path

    if catalog or _fallback_image_path():
        for w_start in range(0, len(pages), IMAGE_SLIDE_CADENCE):
            window = pages[w_start:w_start + IMAGE_SLIDE_CADENCE]
            best_page_idx, best_score, best_path = None, 0, None
            for offset, (page_items, _) in enumerate(window):
                toks = _slide_tokens(title, page_items)
                path, score = _best_image_match(toks, catalog)
                if score > best_score:
                    best_page_idx, best_score, best_path = w_start + offset, score, path
            if best_page_idx is not None:
                image_for_page[best_page_idx] = best_path
                _used_image_paths.add(best_path)
            else:
                fb = _fallback_image_path()
                if fb and fb not in _used_image_paths:
                    image_for_page[w_start] = fb
                    _used_image_paths.add(fb)
                    imagery_warnings.append(
                        f"slide window {w_start // IMAGE_SLIDE_CADENCE + 1}: "
                        f"no keyword match → used fallback image "
                        f"({ONEDRIVE_FALLBACK_FILENAME})"
                    )
                elif fb:
                    imagery_warnings.append(
                        f"slide window {w_start // IMAGE_SLIDE_CADENCE + 1}: "
                        f"no keyword match AND fallback already used — "
                        f"skipped imagery for this window to keep images unique"
                    )

    # ── Build pass ──
    all_slides = []
    image_slide_seq = 0   # global counter (drives S10 every-3rd cadence)

    for p_idx, (page_items, continues) in enumerate(pages):
        t = title if p_idx == 0 else (title + " (continued)")

        if p_idx in image_for_page:
            # Image slide — choose S9 / S10 / S11 variant
            image_slide_seq += 1
            title_fits_1line = (
                _text_width_pt(t, CONTENT_S7_TITLE_PT)
                <= CONTENT_TITLE_WIDTH_PT - SAFETY_PT
            )
            if not title_fits_1line:
                variant = CONTENT_S11_SLIDE_INDEX     # multi-line title → S11
            elif image_slide_seq % S10_CADENCE_AMONG_IMAGES == 0:
                variant = CONTENT_S10_SLIDE_INDEX     # every 3rd image → S10
            else:
                variant = CONTENT_S9_SLIDE_INDEX

            built = _build_image_slide(
                prs, variant, t, page_items, image_for_page[p_idx],
                bullet_mode=bullet_mode,
                first_l1_continues=continues,
                cont_title=title + " (continued)",
            )
            all_slides.extend(built)
        else:
            # Text-only slide
            slide_idx = pick_content_slide_index(t)
            slide = duplicate_slide(prs, slide_idx)
            set_content_title(slide, t)
            set_content_bullets(slide, 1, page_items,
                                bullet_mode=bullet_mode,
                                first_l1_continues=continues)
            all_slides.append(slide)

    return all_slides


def set_content_series(prs, sections, bullet_mode="default"):
    """Build content slides for a series of (title, items) sections so that
    auto-imagery cadence is maintained ACROSS the whole series — every
    `IMAGE_SLIDE_CADENCE` text pages, the best-matching page in that window
    becomes an image slide. Use this instead of calling `set_content_slides`
    per subsection when the subsections belong to the same logical section
    of a deck (so an image is not forced onto every short subsection).

    `sections` is a list of `(title, items)` tuples. Returns the full list
    of slides built across all sections.

    Build-time routing warnings: each section's title is checked for a
    number-keyword ('two', 'three', …, 'ten' or digits 2–10). If the title's
    count matches the L1 bullet count, the section is a STRONG graphic
    candidate. Routing it to plain content (here, instead of S14/S15/S16) is
    almost certainly a missed graphic-routing opportunity, so a line is
    appended to `routing_warnings` and surfaced after the build.
    """
    imagery_warnings.clear()
    catalog = _scan_onedrive_image_catalog()

    # Routing audit — fires on titles where 'two'/'three'/… matches the L1
    # bullet count. The subsection arrived at `set_content_series` (the
    # text/image path) rather than at a graphic helper (S14/S15/S16), so the
    # warning surfaces what was missed during Q2c pre-analysis routing. The
    # subsection may still later promote to a Content+Image variant via
    # cadence or sparse-text rules — the warning is about the GRAPHIC
    # routing decision, not about whether the slide ended up text-only.
    # Don't clear — multiple `set_content_series` calls accumulate across
    # the whole build.
    for section_title, items in sections:
        n = _title_number_keyword(section_title)
        if n is None:
            continue
        l1 = _count_l1_bullets(items)
        if l1 == n:
            routing_warnings.append(
                f"'{section_title}' — title has '{n}' AND content has {l1} "
                f"L1 bullets (matches). Strong candidate for a graphic "
                f"(S14 timeline / S15 target / S16 takeaways). Did not "
                f"route through a graphic helper. Per Q2c rule: graphic "
                f"first, image fallback — verify this is intentional."
            )

    # Body dims from S7 (text-only pagination measure)
    s7_tmpl = prs.slides[CONTENT_S7_SLIDE_INDEX]
    body_w_emu, body_h_emu = _get_ph_dims_from_xml(s7_tmpl, 1)
    body_w_pt = body_w_emu / 12700
    body_h_pt = body_h_emu / 12700
    budget_h_pt = body_h_pt * CONTENT_USABLE_HEIGHT_FRACTION

    # Phase 1 — paginate each section, build a global page list. Each entry:
    #   (section_title, page_items, continues_flag, is_first_page_of_section)
    global_pages = []
    for section_title, items in sections:
        section_pages = _paginate_content_items(
            items, body_w_pt, budget_h_pt, bullet_mode
        )
        for idx, (page_items, continues) in enumerate(section_pages):
            global_pages.append((section_title, page_items, continues, idx == 0))

    # Phase 2 — pick image pages across the entire global page list.
    #
    # 2a) Sparse-text promotion (rule from the brand-ppt polish notes): any
    #     page whose body is sparse enough to leave obvious empty space
    #     ALWAYS becomes an image slide, regardless of cadence. The page's
    #     own tokens drive the image match — fall back to the brand fallback
    #     image when no good keyword match exists.
    # 2b) Cadence-driven imagery for remaining windows: every 4 pages, pick
    #     the best-matching page that hasn't already been promoted. Windows
    #     where a sparse page already grabbed the slot are skipped.
    image_for_page = {}
    have_imagery = bool(catalog or _fallback_image_path())

    if not have_imagery:
        # No-OneDrive mode: same cadence + sparse-promotion structure as the
        # imagery path, but assign the GREY_BOX_PLACEHOLDER sentinel instead
        # of a matched image path. Downstream `_build_image_slide` recognises
        # the sentinel and skips the picture insertion — the template's empty
        # picture placeholder remains in place so PowerPoint renders it as a
        # recognisable "image goes here" box for the user to fill in later.
        for p_idx, (_, page_items, _, _) in enumerate(global_pages):
            if _is_sparse_page(page_items):
                image_for_page[p_idx] = GREY_BOX_PLACEHOLDER
        for w_start in range(0, len(global_pages), IMAGE_SLIDE_CADENCE):
            window_indices = list(range(
                w_start, min(w_start + IMAGE_SLIDE_CADENCE, len(global_pages))
            ))
            if any(i in image_for_page for i in window_indices):
                continue
            image_for_page[window_indices[0]] = GREY_BOX_PLACEHOLDER

    if have_imagery:
        # 2a — sparse promotion
        for p_idx, (sec_title, page_items, _, _) in enumerate(global_pages):
            if not _is_sparse_page(page_items):
                continue
            toks = _slide_tokens(sec_title, page_items)
            path, score = _best_image_match(toks, catalog)
            if not path:
                fb = _fallback_image_path()
                if fb and fb not in _used_image_paths:
                    path = fb
                    imagery_warnings.append(
                        f"sparse page '{sec_title}': no keyword match → used "
                        f"fallback image ({ONEDRIVE_FALLBACK_FILENAME})"
                    )
                elif fb:
                    imagery_warnings.append(
                        f"sparse page '{sec_title}': no keyword match AND "
                        f"fallback already used — skipped imagery to keep "
                        f"images unique"
                    )
            if path:
                image_for_page[p_idx] = path
                _used_image_paths.add(path)

        # 2b — cadence
        for w_start in range(0, len(global_pages), IMAGE_SLIDE_CADENCE):
            window_indices = list(range(
                w_start, min(w_start + IMAGE_SLIDE_CADENCE, len(global_pages))
            ))
            if any(i in image_for_page for i in window_indices):
                continue  # sparse promotion already claimed this window's slot
            best_idx, best_score, best_path = None, 0, None
            for i in window_indices:
                sec_title, page_items, _, _ = global_pages[i]
                toks = _slide_tokens(sec_title, page_items)
                path, score = _best_image_match(toks, catalog)
                if score > best_score:
                    best_idx, best_score, best_path = i, score, path
            if best_idx is not None:
                image_for_page[best_idx] = best_path
                _used_image_paths.add(best_path)
            else:
                fb = _fallback_image_path()
                if fb and fb not in _used_image_paths:
                    image_for_page[w_start] = fb
                    _used_image_paths.add(fb)
                    imagery_warnings.append(
                        f"slide window {w_start // IMAGE_SLIDE_CADENCE + 1}: "
                        f"no keyword match → used fallback image "
                        f"({ONEDRIVE_FALLBACK_FILENAME})"
                    )
                elif fb:
                    imagery_warnings.append(
                        f"slide window {w_start // IMAGE_SLIDE_CADENCE + 1}: "
                        f"no keyword match AND fallback already used — "
                        f"skipped imagery for this window to keep images unique"
                    )

    # Phase 3 — build slides
    all_slides = []
    image_slide_seq = 0

    for p_idx, (sec_title, page_items, continues, is_first_of_section) in enumerate(global_pages):
        # Title: section's title on its first page, "(continued)" on overflow pages
        t = sec_title if is_first_of_section else (sec_title + " (continued)")

        if p_idx in image_for_page:
            image_slide_seq += 1
            title_fits = (_text_width_pt(t, CONTENT_S7_TITLE_PT)
                          <= CONTENT_TITLE_WIDTH_PT - SAFETY_PT)
            if not title_fits:
                variant = CONTENT_S11_SLIDE_INDEX
            elif image_slide_seq % S10_CADENCE_AMONG_IMAGES == 0:
                variant = CONTENT_S10_SLIDE_INDEX
            else:
                variant = CONTENT_S9_SLIDE_INDEX

            built = _build_image_slide(
                prs, variant, t, page_items, image_for_page[p_idx],
                bullet_mode=bullet_mode,
                first_l1_continues=continues,
                cont_title=sec_title + " (continued)",
            )
            all_slides.extend(built)
        else:
            slide_idx = pick_content_slide_index(t)
            slide = duplicate_slide(prs, slide_idx)
            set_content_title(slide, t)
            set_content_bullets(slide, 1, page_items,
                                bullet_mode=bullet_mode,
                                first_l1_continues=continues)
            all_slides.append(slide)

    return all_slides


def _build_image_slide(prs, variant, title, items, image_path,
                       bullet_mode="default", first_l1_continues=False,
                       cont_title=None):
    """Clone S9/S10/S11, set title (Inter Variable Thin), paginate `items` against the
    narrower image-slide body width, insert the picture, and emit S7/S8
    text-only continuation slides for any overflow. Returns the slide list."""
    body_w_pt = {
        CONTENT_S9_SLIDE_INDEX:  IMAGE_S9_BODY_W_PT,
        CONTENT_S10_SLIDE_INDEX: IMAGE_S10_BODY_W_PT,
        CONTENT_S11_SLIDE_INDEX: IMAGE_S11_BODY_W_PT,
    }[variant]
    _, body_h_emu = _get_ph_dims_from_xml(prs.slides[variant], 1)
    body_h_pt = (body_h_emu / 12700) * CONTENT_USABLE_HEIGHT_FRACTION

    sub_pages = _paginate_content_items(items, body_w_pt, body_h_pt, bullet_mode)

    # First sub-page → image slide
    img_slide = duplicate_slide(prs, variant)
    set_content_title(img_slide, title)
    set_content_bullets(img_slide, 1, sub_pages[0][0],
                        bullet_mode=bullet_mode,
                        first_l1_continues=first_l1_continues)
    if image_path == GREY_BOX_PLACEHOLDER:
        # No-OneDrive mode: leave the template's empty picture placeholder in
        # place so PowerPoint renders it as a recognisable "picture goes here"
        # box. Track the slide title for the end-of-build summary.
        grey_box_slides.append(title)
    else:
        try:
            insert_picture(img_slide, 12, image_path, mode="fit")
        except Exception as e:
            imagery_warnings.append(
                f"Failed to insert image {image_path!r}: {type(e).__name__}: {e}"
            )
    slides = [img_slide]

    # Remaining sub-pages → S7/S8 text continuations
    if len(sub_pages) > 1:
        ct = cont_title or (title + " (continued)")
        ct_idx = pick_content_slide_index(ct)
        for page_items, _ in sub_pages[1:]:
            cont_slide = duplicate_slide(prs, ct_idx)
            set_content_title(cont_slide, ct)
            set_content_bullets(cont_slide, 1, page_items,
                                bullet_mode=bullet_mode,
                                first_l1_continues=True)
            slides.append(cont_slide)

    return slides


# ── Phrase-based line breaking (shared by cover title and quote) ─────

BREAK_BEFORE = {
    # Prepositions
    'for', 'of', 'in', 'on', 'at', 'with', 'through', 'across', 'by',
    'to', 'from', 'into', 'over', 'under', 'about', 'against', 'between',
    'during', 'within', 'without', 'beyond', 'beside', 'behind', 'around',
    # Coordinating conjunctions
    'and', 'but', 'or', 'nor', 'yet', 'so',
    # Relative pronouns (start a relative clause)
    'who', 'whom', 'whose', 'which', 'that',
    # Subordinating conjunctions (start a subordinate clause)
    'because', 'since', 'although', 'though', 'while', 'when', 'where',
    'whether', 'if', 'unless', 'until', 'before', 'after', 'as', 'than',
}
NATURAL_PUNCT = (':', ',', ';', '.', '!', '?', '—', '–')
INTRO_PUNCT   = (':', '—', '–', ';')   # only these trigger next-line capitalization
NATURAL_BONUS_PT = 50
SAFETY_PT = 4
# Title widths are measured with Inter Variable Thin via fontTools, but PowerPoint
# renders ~3% wider in practice (subpixel positioning, kerning differences).
# At small placeholder margins this causes the right-edge word to wrap, often
# producing a widow on the next line. Use a larger pad for TITLE layouts.
TITLE_SAFETY_PT = 50
# Quote layouts suffer the same width drift between fontTools and PowerPoint.
# Empirically tuned by rendering the test quote until "runs." no longer
# stranded as a widow in PowerPoint. Larger than the SAFETY_PT=4 default that
# applied before tuning.
QUOTE_SAFETY_PT = 25
LINE_HEIGHT_FACTOR = 1.2   # Inter line-height multiplier used by PowerPoint


def _has_widow(lines):
    """Return True if any line AFTER THE FIRST contains fewer than 2 words.
    A short FIRST line (e.g. `M1:`) is an intentional opener and not a widow
    in the typographic sense — the widow concept refers to a short final line
    of a paragraph. Single-line layouts are never widowed."""
    if len(lines) <= 1:
        return False
    return any(len(line.split()) < 2 for line in lines[1:])


def _layout_phrase_lines(text, ph_w_pt, min_pt=42, max_pt=72, ph_h_pt=None):
    """Find the largest font size (min_pt..max_pt, 2pt steps) and best line split
    for `text` to fit within `ph_w_pt` points of width using Inter Variable Thin metrics.

    If `ph_h_pt` is provided, candidate layouts are also rejected when
    n_lines × size × LINE_HEIGHT_FACTOR exceeds the box height — used for
    placeholders with strict vertical bounds (e.g. section-break titles, only
    ~76pt tall, where a 2-line layout above 30pt would overflow).

    4-phase priority:
      1. Natural 2-line — split at NATURAL_PUNCT or before BREAK_BEFORE word.
      2. Any 2-line — for text with no natural break points.
      3. 3-line fallback.
      4. Greedy fallback at min_pt.

    Capitalizes the first letter of a line only when the previous line ends with
    INTRO_PUNCT. Returns (size_pt, lines).
    """
    # Title layouts use a larger safety margin than the quote/body default,
    # because PowerPoint renders Inter Variable Thin slightly wider than fontTools
    # measures and titles tend to sit right at the placeholder edge.
    fit_w_pt = ph_w_pt - TITLE_SAFETY_PT
    words = text.split()

    def _fits_height(size_pt, n_lines):
        if ph_h_pt is None:
            return True
        return n_lines * size_pt * LINE_HEIGHT_FACTOR <= ph_h_pt

    def _fits(s, size_pt):
        return _text_width_pt(s, size_pt) <= fit_w_pt

    def _is_natural(i):
        return words[i].lower() in BREAK_BEFORE or words[i - 1][-1] in NATURAL_PUNCT

    def _score(lines, size_pt, n_natural):
        return max(_text_width_pt(l, size_pt) for l in lines) - n_natural * NATURAL_BONUS_PT

    def _best_2_natural(size_pt):
        full = ' '.join(words)
        if _fits(full, size_pt):
            return [full]
        best, best_score = None, float('inf')
        for i in range(1, len(words)):
            if not _is_natural(i):
                continue
            l1, l2 = ' '.join(words[:i]), ' '.join(words[i:])
            if _has_widow([l1, l2]):
                continue
            if not _fits(l1, size_pt) or not _fits(l2, size_pt):
                continue
            sc = _score([l1, l2], size_pt, 1)
            if sc < best_score:
                best_score, best = sc, [l1, l2]
        return best

    def _best_2(size_pt):
        full = ' '.join(words)
        if _fits(full, size_pt):
            return [full]
        best, best_score = None, float('inf')
        for i in range(1, len(words)):
            l1, l2 = ' '.join(words[:i]), ' '.join(words[i:])
            if _has_widow([l1, l2]):
                continue
            if not _fits(l1, size_pt) or not _fits(l2, size_pt):
                continue
            sc = _score([l1, l2], size_pt, int(_is_natural(i)))
            if sc < best_score:
                best_score, best = sc, [l1, l2]
        return best

    def _best_3(size_pt):
        result = _best_2(size_pt)
        if result:
            return result
        best, best_score = None, float('inf')
        for i in range(1, len(words)):
            for j in range(i + 1, len(words)):
                l1, l2, l3 = (' '.join(words[:i]), ' '.join(words[i:j]),
                               ' '.join(words[j:]))
                if _has_widow([l1, l2, l3]):
                    continue
                if not all(_fits(l, size_pt) for l in (l1, l2, l3)):
                    continue
                nb = int(_is_natural(i)) + int(_is_natural(j))
                sc = _score([l1, l2, l3], size_pt, nb)
                if sc < best_score:
                    best_score, best = sc, [l1, l2, l3]
        return best

    chosen_lines, chosen_size = None, None

    for size_pt in range(max_pt, min_pt - 1, -2):           # Phase 1 — natural 2-line
        r = _best_2_natural(size_pt)
        if r is not None and _fits_height(size_pt, len(r)):
            chosen_lines, chosen_size = r, size_pt
            break

    if chosen_lines is None:
        for size_pt in range(max_pt, min_pt - 1, -2):       # Phase 2 — any 2-line
            r = _best_2(size_pt)
            if r is not None and _fits_height(size_pt, len(r)):
                chosen_lines, chosen_size = r, size_pt
                break

    if chosen_lines is None:
        for size_pt in range(max_pt, min_pt - 1, -2):       # Phase 3 — 3-line
            r = _best_3(size_pt)
            if r is not None and _fits_height(size_pt, len(r)):
                chosen_lines, chosen_size = r, size_pt
                break

    if chosen_lines is None:                                 # Phase 4 — greedy fallback
        lines, cur = [], []
        for w in words:
            candidate = ' '.join(cur + [w])
            if cur and not _fits(candidate, min_pt):
                lines.append(' '.join(cur))
                cur = [w]
            else:
                cur.append(w)
        if cur:
            lines.append(' '.join(cur))
        chosen_lines, chosen_size = lines, min_pt

    capped = [chosen_lines[0]]
    for i in range(1, len(chosen_lines)):
        prev, line = chosen_lines[i - 1], chosen_lines[i]
        if prev and prev[-1] in INTRO_PUNCT and line:
            line = line[0].upper() + line[1:]
        capped.append(line)
    return chosen_size, capped


def _write_styled_lines(slide, ph_idx, lines, size_pt, typeface=BRAND_FONT_TITLE):
    """Write each line as its own paragraph in the placeholder, preserving the
    template's pPr and stamping each run with the given size + typeface.
    Verifies the read-back line count matches what was written."""
    ph = next(p for p in slide.placeholders if p.placeholder_format.idx == ph_idx)
    tf = ph.text_frame
    txBody = tf._txBody

    template_pPr = None
    if tf.paragraphs:
        pPr_el = tf.paragraphs[0]._element.find(f'{{{nsmap["a"]}}}pPr')
        if pPr_el is not None:
            template_pPr = copy.deepcopy(pPr_el)

    for p in list(txBody.findall(f'{{{nsmap["a"]}}}p')):
        txBody.remove(p)

    for line in lines:
        p_el = etree.SubElement(txBody, f'{{{nsmap["a"]}}}p')
        if template_pPr is not None:
            p_el.insert(0, copy.deepcopy(template_pPr))
        r_el = etree.SubElement(p_el, f'{{{nsmap["a"]}}}r')
        rPr = etree.SubElement(r_el, f'{{{nsmap["a"]}}}rPr')
        rPr.set('lang', 'en-US')
        rPr.set('sz', str(int(size_pt * 100)))
        rPr.set('dirty', '0')
        latin = etree.SubElement(rPr, f'{{{nsmap["a"]}}}latin')
        latin.set('typeface', typeface)
        t_el = etree.SubElement(r_el, f'{{{nsmap["a"]}}}t')
        t_el.text = smart_quotes(line)

    actual = [p.text for p in tf.paragraphs if p.text.strip()]
    assert len(actual) == len(lines), (
        f"_write_styled_lines mismatch in ph_idx={ph_idx}: "
        f"wrote {len(lines)}, read back {len(actual)}"
    )


# ── Quote layout (large-font-first, multi-line) ──────────────

def _layout_quote_lines(text, ph_w_pt, ph_h_pt, min_pt=32, max_pt=72, max_lines=5):
    """Find the LARGEST font (max_pt → min_pt, 2pt step) that fits `text` into
    the box, allowing more lines as needed (constrained by box height).

    Differs from `_layout_phrase_lines`: quotes maximize font size first, line
    count second. At each size we try n_lines = 2, 3, …, height_max — first all-
    natural splits, then any-split — and accept the first fit.

    Returns (size_pt, lines).
    """
    import itertools

    # Quote layouts pay the same fontTools-vs-PowerPoint width-drift tax as
    # title layouts. QUOTE_SAFETY_PT is empirically tuned by rendering until
    # widow-free in PowerPoint — not a theoretical pad.
    fit_w_pt = ph_w_pt - QUOTE_SAFETY_PT
    fit_h_pt = ph_h_pt
    words = text.split()

    def _is_natural(i):
        return words[i].lower() in BREAK_BEFORE or words[i - 1][-1] in NATURAL_PUNCT

    def _lines_from(splits):
        bounds = (0,) + tuple(splits) + (len(words),)
        return [' '.join(words[bounds[i]:bounds[i + 1]]) for i in range(len(bounds) - 1)]

    SENT_END = ('.', '!', '?')
    BALANCE_PENALTY = 0.5     # weight on (max-min) line-width spread

    def _has_midline_sentence_end(lines):
        for line in lines:
            line_words = line.split()
            for w in line_words[:-1]:                 # all but the last word
                if w.endswith(SENT_END):
                    return True
        return False

    def _score(splits, sz):
        lines = _lines_from(splits)
        if _has_widow(lines):
            return None                                # invalid — single-word line
        if _has_midline_sentence_end(lines):
            return None                                # invalid — splits a sentence mid-line
        widths = [_text_width_pt(l, sz) for l in lines]
        if max(widths) > fit_w_pt:
            return None
        n_nat = sum(1 for s in splits if _is_natural(s))
        spread = max(widths) - min(widths)             # penalize unbalanced layouts
        return max(widths) + spread * BALANCE_PENALTY - n_nat * NATURAL_BONUS_PT

    def _best(n, sz, natural_only):
        best, best_sc = None, float('inf')
        for splits in itertools.combinations(range(1, len(words)), n - 1):
            if natural_only and not all(_is_natural(s) for s in splits):
                continue
            sc = _score(splits, sz)
            if sc is not None and sc < best_sc:
                best_sc, best = sc, splits
        return best

    def _capitalize(lines):
        out = [lines[0]]
        for i in range(1, len(lines)):
            prev, line = lines[i - 1], lines[i]
            if prev and prev[-1] in INTRO_PUNCT and line:
                line = line[0].upper() + line[1:]
            out.append(line)
        return out

    for size_pt in range(max_pt, min_pt - 1, -2):
        n_max_height = max(2, int(fit_h_pt / (size_pt * LINE_HEIGHT_FACTOR)))
        n_max = min(max_lines, n_max_height, len(words))

        # Try smallest n first at this font size (prefer fewer lines at same font)
        for n in range(2, n_max + 1):
            splits = _best(n, size_pt, natural_only=True)
            if splits is None:
                splits = _best(n, size_pt, natural_only=False)
            if splits is not None:
                return size_pt, _capitalize(_lines_from(splits))

    # Final fallback: greedy at min_pt, may overflow — caller should warn
    lines, cur = [], []
    for w in words:
        candidate = ' '.join(cur + [w])
        if cur and _text_width_pt(candidate, min_pt) > fit_w_pt:
            lines.append(' '.join(cur))
            cur = [w]
        else:
            cur.append(w)
    if cur:
        lines.append(' '.join(cur))
    return min_pt, _capitalize(lines)


# ── Cover title (S1) ──────────────────────────────────────────

def set_cover_title(slide, title, min_pt=42, max_pt=72, ph_idx=0):
    """S1 cover slide title polish. Inter Variable Thin, phrase-based line breaks,
    fontTools-measured line widths. Returns (size_pt, lines)."""
    ph_w_emu, _ = _get_ph_dims_from_xml(slide, ph_idx)
    ph_w_pt = ph_w_emu / 12700
    size_pt, lines = _layout_phrase_lines(title, ph_w_pt, min_pt, max_pt)
    _write_styled_lines(slide, ph_idx, lines, size_pt)
    return size_pt, lines


# ── Quote (S6) ────────────────────────────────────────────────

def set_quote(slide, quote_text, attribution, min_pt=32, max_pt=72,
              quote_ph_idx=0, attr_ph_idx=11):
    """S6 quote slide. The opening and closing curly quote marks are positioned
    by the template as separate decorative shapes — DO NOT include quote marks
    in `quote_text`. Uses `_layout_quote_lines` (large-font-first) so the text
    fills the placeholder box vertically as well as horizontally. Writes the
    attribution prefixed with an em-dash to ph_idx=11.
    Returns (size_pt, lines)."""
    text = quote_text.strip()

    ph_w_emu, ph_h_emu = _get_ph_dims_from_xml(slide, quote_ph_idx)
    ph_w_pt = ph_w_emu / 12700
    ph_h_pt = ph_h_emu / 12700
    size_pt, lines = _layout_quote_lines(text, ph_w_pt, ph_h_pt, min_pt, max_pt)
    _write_styled_lines(slide, quote_ph_idx, lines, size_pt)

    attr = attribution.lstrip("—–- ").strip()
    set_ph(slide, attr_ph_idx, "— " + attr)

    return size_pt, lines


# ── Section break (S3 light / S4 dark) ────────────────────────

SECTION_SUPPORT_DEFAULT_TEXT = "Optional text if you need to elaborate about this section"
SECTION_SUPPORT_SHAPE_NAME   = "Text Placeholder 6"


def _find_section_support_shape(slide):
    """Locate the supporting-text shape on a section-break slide. S3 has it as
    a real placeholder (ph_idx=11); S4 has it demoted to a regular shape with
    name 'Text Placeholder 6'. Returns the shape or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 11:
            return ph
    for shape in slide.shapes:
        if shape.name == SECTION_SUPPORT_SHAPE_NAME:
            return shape
    return None


def _set_shape_text(shape, text):
    """Replace all text in a shape with `text`, preserving the first run's
    formatting and clearing extra runs/paragraphs."""
    tf = shape.text_frame
    text = smart_quotes(text)
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
    else:
        tf.paragraphs[0].text = text
    while len(tf.paragraphs) > 1:
        tf.paragraphs[-1]._element.getparent().remove(tf.paragraphs[-1]._element)


def set_section_break(slide, title, supporting_text=None, ph_idx=0):
    """S3 (light) or S4 (dark) section break. Title polished with Inter Variable Thin
    and phrase-aware breaks, constrained to fit the title placeholder's height
    (so a 2-line title at 42pt — which would overflow — is rejected in favour
    of a smaller font).

    Supporting text:
      - If `supporting_text` is given and non-empty: populate the supporting
        shape (ph_idx=11 on S3, named 'Text Placeholder 6' on S4).
      - Otherwise: DELETE the supporting shape entirely. Setting it to an
        empty string leaves the default 'Optional text…' showing through.

    Returns (size_pt, lines).
    """
    ph_w_emu, ph_h_emu = _get_ph_dims_from_xml(slide, ph_idx)
    ph_w_pt = ph_w_emu / 12700
    ph_h_pt = ph_h_emu / 12700
    size_pt, lines = _layout_phrase_lines(
        title, ph_w_pt, min_pt=24, max_pt=42, ph_h_pt=ph_h_pt
    )
    _write_styled_lines(slide, ph_idx, lines, size_pt)

    support = _find_section_support_shape(slide)
    if support is not None:
        if supporting_text and supporting_text.strip():
            _set_shape_text(support, supporting_text.strip())
        else:
            support._element.getparent().remove(support._element)

    return size_pt, lines


# ── Auto-imagery (S9 / S10 / S11 from M1 OneDrive) ────────────

# Content + Image slide indices in the M1 2026 template (19 slides total).
CONTENT_S9_SLIDE_INDEX  = 8     # Image right + 1-line title  (default image slide)
CONTENT_S10_SLIDE_INDEX = 9     # Image left + 1-line title   (every 3rd image slide)
CONTENT_S11_SLIDE_INDEX = 10    # Image right + multi-line title (36pt — S8 counterpart)

# Cadence
IMAGE_SLIDE_CADENCE      = 4    # 1 image per 4 content slides (1-in-3-to-4 target)
S10_CADENCE_AMONG_IMAGES = 3    # every 3rd image slide uses S10 (image-left)

# Sparse-text thresholds — pages whose content fills only a small share of the
# body placeholder are promoted to an image variant (S9/S10/S11) even when the
# cadence wouldn't otherwise schedule one. A sparse page is one where:
#   • L1 bullet count ≤ SPARSE_L1_MAX, AND
#   • total word count ≤ SPARSE_WORD_MAX, AND
#   • no L3-or-deeper nested bullets (deep nesting signals density).
# These thresholds were chosen from the screenshot annotation ("3 short bullets
# leave obvious empty space underneath that an image can fill") — they're
# tunable but should remain conservative so we don't promote slides that have
# already filled their vertical real estate.
SPARSE_L1_MAX   = 3
SPARSE_WORD_MAX = 60


def _is_sparse_page(page_items):
    """Return True if a page's `items` list (list of `(level, text)` tuples) is
    sparse enough that a content+image variant (S9/S10/S11) will fit cleanly.

    A page qualifies when it has at most 3 L1 bullets, at most ~60 total words,
    and no L3-or-deeper nested bullets. Sparse pages auto-promote to an image
    variant in `set_content_series` regardless of the cadence window."""
    l1_count = 0
    word_count = 0
    for level, text in page_items:
        if level == 0:
            l1_count += 1
        if level >= 2:
            return False
        word_count += len(str(text).split())
    return l1_count <= SPARSE_L1_MAX and word_count <= SPARSE_WORD_MAX


# Number-keyword → integer mapping. Used to flag plain-content subsections
# whose TITLE has a count word AND whose L1 bullet count matches the count.
# Those subsections are strong graphic candidates (S14 timeline / S15 target /
# S16 takeaways) and being routed to plain S7/S8 is almost certainly a missed
# graphic-routing opportunity.
_TITLE_NUMBER_WORDS = {
    "two": 2, "three": 3, "four": 4, "five": 5, "six": 6,
    "seven": 7, "eight": 8, "nine": 9, "ten": 10,
}


def _title_number_keyword(title):
    """Return the integer count implied by a count-word in `title`, or None.

    Recognized: 'two'/'2', 'three'/'3', … 'ten'/'10'. Case-insensitive,
    matches whole words only (so 'tworld' doesn't match 'two'). The first
    match wins. Digit forms (e.g. '3 stages') are also recognized.
    """
    if not title:
        return None
    lower = title.lower()
    for word, n in _TITLE_NUMBER_WORDS.items():
        if re.search(rf"\b{word}\b", lower):
            return n
    m = re.search(r"\b([2-9]|10)\b", lower)
    if m:
        return int(m.group(1))
    return None


def _count_l1_bullets(items):
    """Return the number of top-level (`level == 0`) bullets in an items list."""
    return sum(1 for level, _ in items if level == 0)


# Build-time warnings about subsections that should probably have been routed
# to a graphic (S14/S15/S16) but went to plain S7/S8 instead. Surfaced
# alongside `imagery_warnings` at the end of each build.
routing_warnings = []

# Image-slide body widths (used to paginate image-slide content)
IMAGE_S9_BODY_W_PT  = 420.5     # S9 body: 5.84"
IMAGE_S10_BODY_W_PT = 439.2     # S10 body: 6.10"
IMAGE_S11_BODY_W_PT = IMAGE_S9_BODY_W_PT

# OneDrive image catalog
ONEDRIVE_IMAGES_DIR = os.path.expanduser(
    "~/Library/CloudStorage/OneDrive-M1Finance/_Creative Assets/"
    "_Creative Asset Hub/Images-ArticlesStockImages/in/img-OrigStockDownloads"
)
ONEDRIVE_FALLBACK_FILENAME = "HowWeDoIt-AI-midjourney.png"
IMAGE_EXTS = (".jpg", ".jpeg", ".png")

# Stop-words stripped from slide tokens and filename tokens during matching.
_IMG_STOPWORDS = {
    "a", "an", "the", "and", "or", "but", "of", "in", "on", "at", "to",
    "for", "with", "by", "from", "into", "as", "is", "are", "be", "been",
    "this", "that", "these", "those", "you", "your", "our", "we", "us",
    "it", "its", "they", "them", "their", "what", "which", "who",
}

# Tokens stripped from filenames specifically (photo-service slugs, etc.)
_IMG_FILENAME_DROP = {
    "istock", "shutterstock", "unsplash", "pexels", "pixabay", "adobe",
    "midjourney", "ai", "img", "image", "photo", "stock", "orig",
}

# Warnings accumulated during a build. Surface to the user after build finishes.
imagery_warnings = []
_onedrive_catalog_cache = None   # filled on first scan

# Sentinel used in `image_for_page` when the build is in no-OneDrive mode.
# `_build_image_slide` recognises this value and leaves the template's empty
# picture placeholder untouched (PowerPoint renders it as a recognisable
# "picture goes here" box) instead of inserting a real image. Slides built
# this way are tracked in `grey_box_slides` for the end-of-build summary so
# the user knows where to drop their own images.
GREY_BOX_PLACEHOLDER = "__GREY_BOX_PLACEHOLDER__"
grey_box_slides = []

# Image paths already assigned in this build — used to keep auto-imagery
# unique across the whole deck. Each successful match adds the chosen path;
# `_best_image_match` skips any path in this set. Call `reset_image_tracking()`
# at the start of a build (or between independent decks in the same process).
_used_image_paths = set()


def reset_image_tracking():
    """Clear the cross-build set of already-used image paths and the
    grey-box tracker. Call at the start of a build when multiple builds
    run in the same Python process (e.g., regression suites)."""
    _used_image_paths.clear()
    grey_box_slides.clear()


def _filename_tokens(name):
    """Tokenize a filename into a set of lowercase keyword tokens.

      'Budgeting-familyKitchenComputer-iStock-1492289682.jpg'
        → {'budgeting', 'family', 'kitchen', 'computer'}

    Splits on `-`, `_`, `.`, and camelCase boundaries. Drops file extension,
    stock-service slugs (`iStock`, `shutterstock`, `unsplash`, ...), pure-digit
    IDs, and short hash-like slugs (random photographer suffixes)."""
    base = os.path.splitext(name)[0]
    # Split on non-alphanumeric AND insert spaces at camelCase boundaries
    spaced = re.sub(r"(?<=[a-z])(?=[A-Z])", " ", base)
    raw = re.split(r"[^A-Za-z0-9]+", spaced)
    tokens = set()
    for t in raw:
        t = t.lower()
        if not t:
            continue
        if t.isdigit():
            continue
        if t in _IMG_FILENAME_DROP or t in _IMG_STOPWORDS:
            continue
        # Drop short hash-like slugs (e.g. 'cEukkv42O40', 'S2W3Q' from
        # photographer suffixes) — these tend to be long and vowel-poor.
        if len(t) >= 6 and sum(c in "aeiouy" for c in t) <= 1:
            continue
        tokens.add(t)
    return tokens


def _scan_onedrive_image_catalog():
    """Return a dict mapping image absolute paths to their filename-token set.
    Returns an empty dict if the OneDrive folder isn't accessible (e.g. user
    is offline, on a different machine, or has the folder unmounted).
    Cached at module level so multiple slides per build re-use the scan."""
    global _onedrive_catalog_cache
    if _onedrive_catalog_cache is not None:
        return _onedrive_catalog_cache

    if os.environ.get("M1_FORCE_NO_ONEDRIVE") == "1":
        imagery_warnings.append(
            "OneDrive image catalog skipped — M1_FORCE_NO_ONEDRIVE=1. "
            "Deck will build text-only."
        )
        _onedrive_catalog_cache = {}
        return _onedrive_catalog_cache

    if not os.path.isdir(ONEDRIVE_IMAGES_DIR):
        imagery_warnings.append(
            f"OneDrive image folder not found at {ONEDRIVE_IMAGES_DIR}. "
            "Deck will build text-only. Check that OneDrive is mounted."
        )
        _onedrive_catalog_cache = {}
        return _onedrive_catalog_cache

    catalog = {}
    for name in sorted(os.listdir(ONEDRIVE_IMAGES_DIR)):
        if not name.lower().endswith(IMAGE_EXTS):
            continue
        path = os.path.join(ONEDRIVE_IMAGES_DIR, name)
        catalog[path] = _filename_tokens(name)
    _onedrive_catalog_cache = catalog
    return catalog


def _slide_tokens(title, items):
    """Return a set of lowercase content tokens from `title` + the L1 subheads
    and L3 bullets of `items`. Drops stopwords and pure-number tokens.
    Used as the left-hand side of the token-overlap match against the
    OneDrive catalog."""
    text_parts = [title or ""]
    for level, text in items:
        if level in (0, 2):   # L1 subheads + L3 bullets are most visualizable
            text_parts.append(text)
    blob = " ".join(text_parts)

    raw = re.split(r"[^A-Za-z0-9]+", blob)
    tokens = set()
    for t in raw:
        t = t.lower()
        if not t or t.isdigit():
            continue
        if t in _IMG_STOPWORDS:
            continue
        tokens.add(t)
    return tokens


def _best_image_match(slide_tokens, catalog, exclude=None):
    """Return (path, score) for the image in `catalog` with the highest token
    overlap against `slide_tokens`, EXCLUDING any path in `exclude` (default
    is the module-level `_used_image_paths` set). Score = |intersection|.
    Returns (None, 0) if no unused overlap > 0. Ties broken alphabetically
    by path for determinism."""
    if exclude is None:
        exclude = _used_image_paths
    best_path, best_score = None, 0
    for path in sorted(catalog.keys()):
        if path in exclude:
            continue
        score = len(slide_tokens & catalog[path])
        if score > best_score:
            best_path, best_score = path, score
    return best_path, best_score


def _fallback_image_path():
    """Return the absolute path to the designated generic fallback image,
    or None if OneDrive is unavailable or the fallback file is missing.
    Honours the `M1_FORCE_NO_ONEDRIVE` test hook so we can verify the truly
    text-only fallback path."""
    if os.environ.get("M1_FORCE_NO_ONEDRIVE") == "1":
        return None
    path = os.path.join(ONEDRIVE_IMAGES_DIR, ONEDRIVE_FALLBACK_FILENAME)
    return path if os.path.exists(path) else None


# ── Graphic-slide populators (S14 timeline, S15 target, S16 takeaways) ─

def _find_shape_by_name(slide_or_group, name):
    """Return the first shape with the given `name` on a slide OR inside a
    group, or None. Searches `.shapes` so it works for both."""
    for shape in slide_or_group.shapes:
        if shape.name == name:
            return shape
    return None


def _set_two_para_text(shape, line1, line2):
    """Write `line1` into the FIRST paragraph and `line2` into the SECOND
    paragraph of a textbox, preserving each paragraph's existing run-level
    styling (so the bold subhead stays bold, the regular takeaway stays
    regular). Designed for the 'Subhead | Takeaway' template textboxes on
    S15 and S16."""
    tf = shape.text_frame
    if len(tf.paragraphs) >= 1 and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = smart_quotes(line1)
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
    if len(tf.paragraphs) >= 2 and tf.paragraphs[1].runs:
        tf.paragraphs[1].runs[0].text = smart_quotes(line2)
        for r in tf.paragraphs[1].runs[1:]:
            r.text = ""


def _keep_first_n_paragraphs(shape, n):
    """Drop all paragraphs after the first `n` in a textbox. Used to
    collapse S16's TextBox 26 (which holds two subhead-pairs) down to one."""
    tf = shape.text_frame
    txBody = tf._txBody
    paras = txBody.findall(f"{{{nsmap['a']}}}p")
    for p in paras[n:]:
        txBody.remove(p)


def _remove_shape(slide, name):
    """Delete a shape from the slide by name (if present)."""
    shape = _find_shape_by_name(slide, name)
    if shape is not None:
        shape._element.getparent().remove(shape._element)


def _clone_shape_after(slide, src_shape, new_name=None):
    """Deep-copy `src_shape` and insert the clone directly after the source in
    the slide's spTree. Returns the new shape (resolved via element identity).
    Optionally renames the clone so `_find_shape_by_name` can disambiguate.

    Assigns a UNIQUE `cNvPr id` to the clone — without this, the clone keeps
    the source's id and PowerPoint flags the slide as needing repair (each
    shape on a slide must have a unique non-visual id)."""
    import copy as _copy
    new_el = _copy.deepcopy(src_shape._element)
    nv = new_el.find(f"{{{nsmap['p']}}}nvSpPr/{{{nsmap['p']}}}cNvPr")
    if nv is not None:
        # Assign a unique id = max existing id on slide + 1
        existing_ids = []
        for cnv in slide.shapes._spTree.iter(f"{{{nsmap['p']}}}cNvPr"):
            try:
                existing_ids.append(int(cnv.get("id", "0")))
            except (TypeError, ValueError):
                pass
        nv.set("id", str(max(existing_ids, default=0) + 1))
        if new_name is not None:
            nv.set("name", new_name)
    src_shape._element.addnext(new_el)
    for sh in slide.shapes:
        if sh._element is new_el:
            return sh
    return None


def _set_shape_xywh(shape, left=None, top=None, width=None, height=None):
    """Set absolute position+size on a shape, leaving unspecified dims as-is."""
    if left   is not None: shape.left   = int(left)
    if top    is not None: shape.top    = int(top)
    if width  is not None: shape.width  = int(width)
    if height is not None: shape.height = int(height)


def _as_label_text_pair(item):
    """Normalize a graphic-helper input item into a `(label, text)` pair.
    Accepts a plain string (label, empty text) or a `(label, text)` tuple."""
    if isinstance(item, (tuple, list)):
        return (item[0], item[1] if len(item) > 1 else "")
    return (item, "")


# ── S14 geometry constants (extracted directly from the template) ──
_S14_SLIDE_WIDTH    = 12192000   # EMU
_S14_OVAL_TOP       = 2014722
_S14_OVAL_SIZE      = 2093132    # ovals are square (w == h)
_S14_OVAL_LEFTS     = [1270884, 3203016, 5135149, 7067281, 8999414]   # 5-position lefts
_S14_GROUP_TOP      = 2657458
_S14_GROUP_LEFTS    = [1471422, 3407125, 5347600, 7283303, 9216635]
_S14_MARKER_TOP     = 4336137    # year/date marker textbox top
_S14_MARKER_LEFTS   = [1988880, 3940656, 5892432, 7844207, 9795982]
_S14_CONNECTOR_TOP  = 4500422    # dotted line y
_S14_CONNECTOR_LEFT = 2135704    # dotted line left edge (centered under marker 0)
_S14_MARKER_OFFSET_FROM_OVAL = _S14_MARKER_LEFTS[0] - _S14_OVAL_LEFTS[0]  # 717996
_S14_GROUP_OFFSET_FROM_OVAL  = _S14_GROUP_LEFTS[0]  - _S14_OVAL_LEFTS[0]  # 200538
_S14_CONNECTOR_OFFSET_FROM_OVAL = _S14_CONNECTOR_LEFT - _S14_OVAL_LEFTS[0]  # 864820


def _as_step_tuple(item):
    """Normalize an S14 step entry into `(label, body, marker)`.
    Accepts plain string, 2-tuple, or 3-tuple. Marker is None when absent."""
    if isinstance(item, str):
        return (item, "", None)
    label  = item[0]
    body   = item[1] if len(item) > 1 else ""
    marker = item[2] if len(item) > 2 and item[2] else None
    return (label, body, marker)


def set_s14_timeline(slide, title, steps, takeaway=None):
    """Populate S14 timeline graphic: title + per-step ovals + optional
    year/date markers + dotted connector line + bottom takeaway.

    `steps` — up to 5 entries. Each entry is one of:
        - `"Inform"`                                       (label only)
        - `("Inform", "Answers your questions")`           (label + body)
        - `("Inform", "Answers your questions", "2022")`   (label + body + marker)

    Step markers are timeline dates/years that appear BELOW the dotted line
    underneath each oval. If a step has no marker, that step's marker textbox
    is removed. If NO step has a marker, the dotted connector is removed too.

    Unused ovals/groups (steps[len(steps):5]) are also removed, and the
    remaining cluster of ovals + groups + markers + connector is recentered
    horizontally on the slide. The bottom takeaway is populated when supplied
    and removed otherwise.
    """
    set_content_title(slide, title)
    items = [_as_step_tuple(s) for s in steps[:5]]
    n = len(items)

    step_label_names = ["TextBox 49", "TextBox 50", "TextBox 51", "TextBox 52", "TextBox 53"]
    group_names      = ["Group 32",   "Group 34",   "Group 37",   "Group 40",   "Group 43"]
    oval_names       = ["Oval 20",    "Oval 21",    "Oval 23",    "Oval 24",    "Oval 25"]
    group_text_names = [
        ("TextBox 30", "TextBox 31"),
        ("TextBox 35", "TextBox 36"),
        ("TextBox 38", "TextBox 39"),
        ("TextBox 41", "TextBox 42"),
        ("TextBox 44", "TextBox 45"),
    ]

    # Phase 1 — fill in-oval header+body. Marker goes to the step-label
    # textbox only if supplied; otherwise that step-label is deleted below.
    for i, (label, body, marker) in enumerate(items):
        grp = _find_shape_by_name(slide, group_names[i])
        if grp is not None:
            hdr_name, body_name = group_text_names[i]
            hdr_shape = _find_shape_by_name(grp, hdr_name)
            bdy_shape = _find_shape_by_name(grp, body_name)
            if hdr_shape is not None:
                _set_shape_text(hdr_shape, label)
            if bdy_shape is not None:
                _set_shape_text(bdy_shape, body)
        lbl = _find_shape_by_name(slide, step_label_names[i])
        if lbl is not None:
            if marker:
                _set_shape_text(lbl, marker)
            else:
                _remove_shape(slide, step_label_names[i])

    # Phase 2 — remove unused step slots
    for i in range(n, 5):
        for nm in (step_label_names[i], group_names[i], oval_names[i]):
            _remove_shape(slide, nm)

    # Phase 3 — dotted connector: delete if no step in [0..n) has a marker;
    # otherwise trim its right edge to align with the last visible marker
    any_marker = any(m for _, _, m in items)
    if not any_marker:
        _remove_shape(slide, "Straight Connector 47")

    # Phase 4 — recenter the cluster horizontally on the slide. Compute new
    # left positions for the n kept ovals (preserving the original overlap
    # spacing) so the cluster's midpoint aligns with the slide's midpoint.
    if n < 5:
        # Original n-oval cluster span (left edge of oval 0 → right edge of oval n-1)
        orig_left  = _S14_OVAL_LEFTS[0]
        orig_right = _S14_OVAL_LEFTS[n - 1] + _S14_OVAL_SIZE
        cluster_w  = orig_right - orig_left
        new_left_0 = (_S14_SLIDE_WIDTH - cluster_w) // 2
        shift      = new_left_0 - orig_left

        for i in range(n):
            shape = _find_shape_by_name(slide, oval_names[i])
            if shape is not None:
                shape.left = shape.left + shift
            shape = _find_shape_by_name(slide, group_names[i])
            if shape is not None:
                shape.left = shape.left + shift
            shape = _find_shape_by_name(slide, step_label_names[i])
            if shape is not None:
                shape.left = shape.left + shift

        # Connector — only if any marker exists; trim width to span from
        # marker[0] center to marker[n-1] center (matches template behavior
        # where the line spans across the visible markers).
        if any_marker:
            conn = _find_shape_by_name(slide, "Straight Connector 47")
            if conn is not None:
                new_conn_left  = _S14_CONNECTOR_LEFT + shift
                new_conn_right = (_S14_MARKER_LEFTS[n - 1] + shift) + 657139  # marker w
                conn.left  = new_conn_left
                conn.width = new_conn_right - new_conn_left

    # Phase 5 — bottom takeaway summary
    if takeaway:
        tk = _find_shape_by_name(slide, "TextBox 54")
        if tk is not None:
            _set_shape_text(tk, takeaway)
    else:
        _remove_shape(slide, "TextBox 54")


# ── S15 geometry constants (extracted directly from the template) ──
_S15_CENTER_X       = 3003010   # shared x-center of all ovals
_S15_BOTTOM_Y       = 5956898   # shared bottom-y of all ovals
_S15_ORIG_SIZES     = [1465617, 2093132, 3282632, 4222498]  # inner → outer
_S15_PAIR_FIRST_TOP = 1752516   # top pair textbox (outermost = position n-1)
_S15_PAIR_LAST_TOP  = 4859643   # bottom pair textbox (innermost = position 0)
_S15_PAIR_HEIGHT    = 755060
_S15_MARKER_DY      = -66364    # marker stripe y-offset relative to pair textbox
_S15_LABEL_HEIGHT   = 230832
_S15_LABEL_WIDTH_RATIO = 0.8    # ring label width as fraction of oval width


def _s15_compute_layout(n, enlarge):
    """Compute the full S15 geometry for `n` visible rings (innermost first).

    When `enlarge` is True and n < 4, every kept ring is scaled so the
    outermost-kept ring matches the template's original outermost size. All
    ovals stay concentric (same x-center, same bottom-y). Pair textboxes
    are redistributed evenly between the original top and bottom positions.
    """
    scale = _S15_ORIG_SIZES[3] / _S15_ORIG_SIZES[n - 1] if (enlarge and n < 4) else 1.0

    ovals = []
    for i in range(n):
        size = int(_S15_ORIG_SIZES[i] * scale)
        ovals.append({
            "left": _S15_CENTER_X - size // 2,
            "top":  _S15_BOTTOM_Y - size,
            "size": size,
        })

    # Ring labels: width = 80% of oval, centered horizontally; vertical
    # placement = upper third of each ring's visible sliver (between this
    # ring's top and the next inner ring's top). Innermost ring is centered.
    labels = []
    for i in range(n):
        o = ovals[i]
        lbl_w = int(o["size"] * _S15_LABEL_WIDTH_RATIO)
        if i == 0:
            lbl_center_y = o["top"] + o["size"] // 2
        else:
            inner_top = ovals[i - 1]["top"]
            lbl_center_y = o["top"] + int(0.3 * (inner_top - o["top"]))
        labels.append({
            "left":   _S15_CENTER_X - lbl_w // 2,
            "top":    lbl_center_y - _S15_LABEL_HEIGHT // 2,
            "width":  lbl_w,
            "height": _S15_LABEL_HEIGHT,
        })

    # Pair textboxes — innermost = bottom, outermost = top
    if n == 1:
        pair_tops = [(_S15_PAIR_FIRST_TOP + _S15_PAIR_LAST_TOP) // 2]
    else:
        step = (_S15_PAIR_LAST_TOP - _S15_PAIR_FIRST_TOP) / (n - 1)
        pair_tops = [int(_S15_PAIR_LAST_TOP - i * step) for i in range(n)]

    return {
        "ovals":   ovals,
        "labels":  labels,
        "pairs":   pair_tops,
        "markers": [t + _S15_MARKER_DY for t in pair_tops],
    }


def _get_content_text_lstStyle(prs):
    """Return a deep copy of the S7 'Content Text' layout's body placeholder
    `<a:lstStyle>`. This contains the brand's L1 body text spec (bold, dark
    navy, 18pt, no bullet) — which is NOT defined in any master, only in
    this one layout. Used to bake the L1 styling into shapes moved to other
    slides (e.g., S15) so the styling no longer depends on layout
    inheritance."""
    import copy
    a_ns = nsmap["a"]
    p_ns = nsmap["p"]
    # The template has 3 slide masters; `prs.slide_layouts` only returns
    # layouts from the first master. Iterate every master's layouts to find
    # "Content Text" (which lives in master 2).
    layout = None
    for master in prs.slide_masters:
        for l in master.slide_layouts:
            if l.name.strip() == "Content Text":
                layout = l
                break
        if layout is not None:
            break
    if layout is None:
        return None
    # Body placeholder = <p:sp> with <p:ph type="body"/> or idx="1"
    for sp in layout.shapes:
        ph_el = sp._element.find(f".//{{{p_ns}}}ph")
        if ph_el is None:
            continue
        if ph_el.get("type") == "body" or ph_el.get("idx") == "1":
            lst = sp._element.find(f".//{{{a_ns}}}lstStyle")
            if lst is not None:
                return copy.deepcopy(lst)
    return None


def _add_s15_summary(slide, summary):
    """Add a summary block above the rings using the L1 body text style.

    The L1 spec (bold, dark navy `#152B56`, 18pt, no bullet) lives only in
    the S7 'Content Text' layout's body placeholder `<a:lstStyle>` — not in
    any master. Strategy: build the text on a temp S7 slide so `set_ph`
    inherits styling correctly, deep-copy the populated body placeholder to
    S15, then BAKE the S7 layout's `<a:lstStyle>` into the copied shape so
    the L1 styling becomes local and doesn't depend on layout inheritance.
    """
    import copy
    a_ns = nsmap["a"]
    prs = slide.part.package.presentation_part.presentation

    layout_lst = _get_content_text_lstStyle(prs)

    temp = duplicate_slide(prs, CONTENT_S7_SLIDE_INDEX)
    set_ph(temp, 1, summary)
    body = next((p for p in temp.placeholders if p.placeholder_format.idx == 1), None)
    if body is not None:
        left, top, width, height = body.left, body.top, body.width, body.height
        new_el = copy.deepcopy(body._element)

        # Bake the S7 layout's L1 spec into the shape's own <a:lstStyle> so
        # the styling no longer depends on the destination slide's layout.
        # Note: <p:txBody> lives in the `p` namespace, but its children
        # <a:bodyPr>, <a:lstStyle>, <a:p> are all in the `a` namespace.
        if layout_lst is not None:
            txBody = new_el.find(f".//{{{nsmap['p']}}}txBody")
            if txBody is not None:
                existing = txBody.find(f"{{{a_ns}}}lstStyle")
                if existing is not None:
                    txBody.remove(existing)
                bodyPr = txBody.find(f"{{{a_ns}}}bodyPr")
                insert_at = list(txBody).index(bodyPr) + 1 if bodyPr is not None else 0
                txBody.insert(insert_at, copy.deepcopy(layout_lst))

        slide.shapes._spTree.append(new_el)
        for sh in slide.shapes:
            if sh._element is new_el:
                _set_shape_xywh(sh, left=left, top=top, width=width, height=height)
                break
    delete_slide(prs, len(prs.slides) - 1)


def set_s15_target(slide, title, levels, summary=None, enlarge=None):
    """Populate S15 target graphic: title + 2–4 nested rings + right-side
    subhead/takeaway pairs.

    `levels` — 2 to 4 entries, ordered **INNERMOST → OUTERMOST**.
        - items[0]  → innermost ring (Core slot — darkest, smallest)
        - items[-1] → outermost ring (Desire slot — lightest, largest)
    Each entry is a plain string (label only) or `(label, takeaway)` tuple.

    Layout rules:
        - Right-side pair textboxes are stacked top→bottom matching the
          template's color order: TOP = outermost (lightest marker),
          BOTTOM = innermost (darkest marker). The helper writes each pair
          into the textbox whose marker color matches its ring.
        - For N<4 levels, the OUTERMOST (largest) rings + their pair
          textboxes + their interior labels are removed. This preserves the
          innermost-first semantic (items[0] always lands on Core).
        - For N>4 the call is clamped to the first 4 levels and a routing
          warning is emitted. The template only ships 4 ovals; adding more
          rings requires manual template editing.

    Mapping (inner → outer):
        ring labels: TextBox 49 (Core), TextBox 2 (Need), TextBox 3 (Utility), TextBox 4 (Desire)
        pair boxes:  TextBox 13       , TextBox 11      , TextBox 9          , TextBox 54
        ovals:       Oval 25          , Oval 23         , Oval 21            , Oval 20
    """
    set_content_title(slide, title)

    if len(levels) > 4:
        routing_warnings.append(
            f"S15 target '{title}' has {len(levels)} levels — clamped to 4 "
            f"(template only ships 4 nested ovals). Trim to ≤4 layers or "
            f"split across two slides."
        )
    if len(levels) < 2:
        routing_warnings.append(
            f"S15 target '{title}' has only {len(levels)} level(s) — graphic "
            f"needs at least 2 concepts that build on each other. Consider a "
            f"plain content slide for 1-concept material."
        )

    items = [_as_label_text_pair(l) for l in levels[:4]]
    n = max(1, len(items))

    # Default mode logic:
    #   summary=None  → enlarge=True (default: graphic fills the freed space)
    #   summary='...' → enlarge=False (summary fills freed space instead)
    #   Caller can pass enlarge=True with summary to get "both" (option c).
    if enlarge is None:
        enlarge = (summary is None)

    # Indexed inner → outer
    ring_names   = ["TextBox 49", "TextBox 2",  "TextBox 3",  "TextBox 4"]
    pair_names   = ["TextBox 13", "TextBox 11", "TextBox 9",  "TextBox 54"]
    oval_names   = ["Oval 25",    "Oval 23",    "Oval 21",    "Oval 20"]
    # Per-pair marker stripe to the left of each subhead textbox (color
    # matches the corresponding ring). Inner → outer.
    marker_names = ["Rounded Rectangle 12", "Rounded Rectangle 10",
                    "Rounded Rectangle 8",  "Rounded Rectangle 5"]

    # Phase 1 — fill text in used slots
    for i, (label, takeaway) in enumerate(items):
        ring = _find_shape_by_name(slide, ring_names[i])
        if ring is not None:
            _set_shape_text(ring, label)
        pair = _find_shape_by_name(slide, pair_names[i])
        if pair is not None:
            _set_two_para_text(pair, label, takeaway)

    # Phase 2 — remove unused outermost rings + scaffolding
    for i in range(n, 4):
        for nm in (ring_names[i], pair_names[i], oval_names[i], marker_names[i]):
            _remove_shape(slide, nm)

    # Phase 3 — geometry. Always run the layout pass so ring labels get
    # widened to fit (kills the "Foundation → Foundati / on" wrap). For
    # n<4 with enlarge=True the rings are also scaled up to fill the
    # original 4-ring footprint.
    layout = _s15_compute_layout(n, enlarge=(enlarge and n < 4))
    for i in range(n):
        oval = _find_shape_by_name(slide, oval_names[i])
        if oval is not None:
            o = layout["ovals"][i]
            _set_shape_xywh(oval, left=o["left"], top=o["top"],
                            width=o["size"], height=o["size"])

        ring_label = _find_shape_by_name(slide, ring_names[i])
        if ring_label is not None:
            l = layout["labels"][i]
            # For N=4 we keep the template's original vertical placement of
            # the labels (small tweaks only — width grows so 'Foundation'
            # fits, but top/height stay where the template designer put them).
            # For N<4 we use the full computed layout.
            if n == 4 and not enlarge:
                _set_shape_xywh(ring_label, left=l["left"], width=l["width"])
            else:
                _set_shape_xywh(ring_label, left=l["left"], top=l["top"],
                                width=l["width"], height=l["height"])

        # Pair textboxes and their marker stripes only get redistributed
        # when enlarge=True. In summary mode the kept pairs stay at their
        # template positions so they remain visually aligned with the
        # (unmoved) rings they describe.
        pair = _find_shape_by_name(slide, pair_names[i])
        if pair is not None and enlarge and n < 4:
            _set_shape_xywh(pair, top=layout["pairs"][i])

        marker = _find_shape_by_name(slide, marker_names[i])
        if marker is not None and enlarge and n < 4:
            _set_shape_xywh(marker, top=layout["markers"][i])

    # Phase 4 — optional summary body text in the freed space above the rings
    if summary:
        _add_s15_summary(slide, summary)


# S16 row geometry (extracted directly from the template's textbox positions).
# Each topic row's "full-width" bbox spans from the left edge of the row's
# leftmost textbox to the right edge of Topic 3's rightmost cell (so a single
# wide proof in any row stretches to the same right margin as the 4-col grid).
S16_ROW_LEFT          = 3328550   # EMU — left edge of every row
S16_ROW_RIGHT         = 11092794  # EMU — right edge (matches Topic 3 col 4 right)
S16_ROW_WIDTH         = S16_ROW_RIGHT - S16_ROW_LEFT       # 7764244 EMU ≈ 8.50"
S16_ROW_GAP           = 270000    # EMU between adjacent columns (~0.30")
S16_ROW_HEIGHT_SINGLE = 755060    # EMU — Topic 1/2's natural 2-paragraph height
S16_ROW_TOPS          = [1809659, 2981082, 4158608]  # EMU — top of each row


def _s16_normalize_topic(item):
    """Normalize one entry in a `topics` list to `(label, proofs)`:

    - Plain string                          → (label, [])               (no proofs)
    - `(label, "single takeaway str")`      → (label, [(label, str)])   (1 wide)
    - `(label, [(sub, take), …])`           → (label, [(sub, take), …]) (N cells)
    - `(label, ["sub: take", "sub: take"])` → splits each on first ':'
    """
    if isinstance(item, str):
        return (item, [])
    label, payload = item[0], item[1] if len(item) > 1 else ""
    if isinstance(payload, str):
        # Single wide proof; reuse the label as the subhead (existing behavior)
        return (label, [(label, payload)] if payload else [])
    proofs = []
    for entry in payload:
        if isinstance(entry, (tuple, list)):
            proofs.append((entry[0], entry[1] if len(entry) > 1 else ""))
        else:
            # "Subhead: takeaway" string — split on first colon
            text = str(entry)
            head, _, tail = text.partition(":")
            proofs.append((head.strip(), tail.strip()) if tail else (text.strip(), ""))
    return (label, proofs)


def _s16_lay_out_row(slide, source_shape, proofs, row_top, row_height):
    """Lay out 1–4 proofs across a single S16 topic row.

    - 1 proof  → `source_shape` resized to full row width
    - 2-4 proofs → `source_shape` becomes column 1; clones become columns 2..N

    Each cell is filled with the proof's (subhead, takeaway) via
    `_set_two_para_text` so the bold-subhead / regular-takeaway styling is
    preserved.
    """
    n = max(1, min(4, len(proofs)))
    col_width = (S16_ROW_WIDTH - (n - 1) * S16_ROW_GAP) // n
    _set_shape_xywh(
        source_shape,
        left=S16_ROW_LEFT,
        top=row_top,
        width=col_width,
        height=row_height,
    )
    _keep_first_n_paragraphs(source_shape, 2)
    _set_two_para_text(source_shape, proofs[0][0], proofs[0][1])

    prev = source_shape
    for col in range(1, n):
        left = S16_ROW_LEFT + col * (col_width + S16_ROW_GAP)
        clone = _clone_shape_after(slide, source_shape, new_name=f"{source_shape.name}_col{col+1}")
        if clone is None:
            continue
        _set_shape_xywh(clone, left=left, top=row_top, width=col_width, height=row_height)
        _keep_first_n_paragraphs(clone, 2)
        _set_two_para_text(clone, proofs[col][0], proofs[col][1])
        prev = clone


def set_s16_takeaways(slide, title, topics):
    """Populate S16 takeaways graphic: title + up to 3 topic rectangles, each
    row flex-laying out 1–4 proof points across one horizontal row.

    `topics` — up to 3 entries. Each entry is one of:

        - `"Topic"`                          → topic pill only, no proofs
        - `("Topic", "takeaway sentence")`   → one WIDE proof spanning full row
        - `("Topic", [proof, …])`            → 1–4 proofs as a flex row

    Each `proof` is `(subhead, takeaway)` tuple OR a `"subhead: takeaway"` str.

    Layout per topic row:
        - 1 proof  → single textbox spans the full right-side width
        - 2 proofs → 2 equal columns with a gap
        - 3 proofs → 3 equal columns with gaps
        - 4 proofs → 4 equal columns with gaps

    The template's extra Topic 3 textboxes (TextBox 31/32/33) are removed up
    front; the helper then clones the per-topic source textbox as many times
    as the proof count requires. Each cell preserves the template's bold-
    subhead / regular-takeaway run styling.
    """
    set_content_title(slide, title)

    # Clear out Topic 3's prebuilt 4-col scaffolding — we re-build dynamically
    for nm in ("TextBox 31", "TextBox 32", "TextBox 33"):
        _remove_shape(slide, nm)

    rect_names    = ["Rounded Rectangle 17", "Rounded Rectangle 18", "Rounded Rectangle 22"]
    source_names  = ["TextBox 24", "TextBox 30", "TextBox 26"]

    used = topics[:3]
    for i, raw in enumerate(used):
        label, proofs = _s16_normalize_topic(raw)
        rect = _find_shape_by_name(slide, rect_names[i])
        if rect is not None:
            _set_shape_text(rect, label)
        source = _find_shape_by_name(slide, source_names[i])
        if source is None:
            continue
        if not proofs:
            # Topic pill only — remove the placeholder textbox so the template
            # 'Subhead | Any key takeaways…' scaffolding doesn't leak through
            source._element.getparent().remove(source._element)
            continue
        _s16_lay_out_row(slide, source, proofs, S16_ROW_TOPS[i], S16_ROW_HEIGHT_SINGLE)

    # Remove unused topic-row scaffolding when fewer than 3 topics are supplied
    for i in range(len(used), 3):
        _remove_shape(slide, rect_names[i])
        _remove_shape(slide, source_names[i])


# ── S12 bar chart / S13 pie chart populators ────────────────

def _find_chart(slide):
    """Return the first chart object on the slide, or None."""
    for sh in slide.shapes:
        if sh.has_chart:
            return sh.chart
    return None


def set_s12_bar(slide, title, data, series_name="Value"):
    """Populate S12 bar chart with real outline data.

    `data` is a list of `(category_label, numeric_value)` pairs, e.g.
        [("Q3 2019", 0), ("Q2 2021", 5), ("Q1 2023", 8), ("Q4 2025", 12)]

    Replaces the template's placeholder Q1–Q4 / Series 1–3 data with a
    single series carrying the real outline values. Hides the template's
    'Chart Title' placeholder."""
    from pptx.chart.data import CategoryChartData
    set_content_title(slide, title)
    chart = _find_chart(slide)
    if chart is None:
        return
    chart_data = CategoryChartData()
    chart_data.categories = [label for label, _ in data]
    chart_data.add_series(series_name, tuple(value for _, value in data))
    chart.replace_data(chart_data)
    chart.has_title = False


def _flip_s13_dlbls_to_percent(chart, num_active_slices):
    """Experiment v2: chart-level flag flip (safe per prior test) PLUS,
    for each active per-point <c:dLbl> (idx 0..num_active_slices-1) inside
    the series-level <c:dLbls>, remove its <c:tx> rich-text override and
    flip its own show flags to percent-only. Leave orphan per-point dLbls
    (idx >= num_active_slices) completely untouched. Leave the series-level
    <c:dLbls> show flags untouched."""
    c_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    c = f"{{{c_ns}}}"
    cs = chart._chartSpace

    # 1. Chart-level: flip show flags (known safe from prior test).
    doughnut = cs.find(f".//{c}doughnutChart")
    if doughnut is None:
        return
    chart_lvl = doughnut.find(f"{c}dLbls")
    if chart_lvl is not None:
        for tag, val in (("showVal", "0"), ("showCatName", "0"), ("showPercent", "1")):
            el = chart_lvl.find(f"{c}{tag}")
            if el is not None:
                el.set("val", val)

    # 2. Per-point dLbls: only for active indices, remove <c:tx> and flip
    # show flags. Leave orphan dLbls intact.
    ser_lvl = doughnut.find(f"{c}ser/{c}dLbls")
    if ser_lvl is None:
        return
    for dLbl in ser_lvl.findall(f"{c}dLbl"):
        idx_el = dLbl.find(f"{c}idx")
        if idx_el is None:
            continue
        try:
            idx = int(idx_el.get("val"))
        except (TypeError, ValueError):
            continue
        if idx >= num_active_slices:
            continue
        tx = dLbl.find(f"{c}tx")
        if tx is not None:
            dLbl.remove(tx)
        # Replace per-point <c:txPr> with a uniform white sz=1800 block.
        old_txPr = dLbl.find(f"{c}txPr")
        if old_txPr is not None:
            dLbl.remove(old_txPr)
        a_ns = nsmap["a"]
        new_txPr = etree.fromstring(
            f'<c:txPr xmlns:c="{c_ns}" xmlns:a="{a_ns}">'
            f'<a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:pPr>'
            f'<a:defRPr sz="1800" b="0">'
            f'<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
            f'<a:latin typeface="+mn-lt"/>'
            f'</a:defRPr>'
            f'</a:pPr><a:endParaRPr lang="en-US"/></a:p>'
            f'</c:txPr>'
        )
        # Schema order: txPr goes after spPr, before showLegendKey.
        anchor = dLbl.find(f"{c}showLegendKey")
        if anchor is not None:
            anchor.addprevious(new_txPr)
        else:
            dLbl.append(new_txPr)
        for tag, val in (("showVal", "0"), ("showCatName", "0"), ("showPercent", "1")):
            el = dLbl.find(f"{c}{tag}")
            if el is not None:
                el.set("val", val)


# Pair-wise mapping for S13 components: (label textbox, marker rectangle)
# Indexed 0..5 by component number (1..6 in the template).
_S13_COMPONENT_SHAPES = [
    ("TextBox 19", "Rounded Rectangle 16"),
    ("TextBox 20", "Rounded Rectangle 15"),
    ("TextBox 21", "Rounded Rectangle 29"),
    ("TextBox 22", "Rounded Rectangle 14"),
    ("TextBox 24", "Rounded Rectangle 23"),
    ("TextBox 28", "Rounded Rectangle 27"),
]


def set_s13_pie(slide, title, slices, series_name="Allocation"):
    """Populate S13 pie chart with real outline data.

    `slices` is a list of `(label, percent)` pairs, e.g.
        [("Invest", 68), ("High-Yield Cash", 18), ("Borrow", 14)]

    Replaces the template's 6-component placeholder with the supplied
    slices and removes the unused per-component label textboxes AND their
    color marker rectangles. Hides chart title and legend so no template
    placeholder text leaks."""
    from pptx.chart.data import CategoryChartData
    set_content_title(slide, title)
    chart = _find_chart(slide)
    if chart is None:
        return
    chart_data = CategoryChartData()
    chart_data.categories = [label for label, _ in slices]
    chart_data.add_series(series_name, tuple(pct for _, pct in slices))
    chart.replace_data(chart_data)
    chart.has_title = False
    chart.has_legend = False
    _flip_s13_dlbls_to_percent(chart, num_active_slices=len(slices))

    # Populate kept components, delete unused ones (both label textbox + marker)
    for i, (label, pct) in enumerate(slices[:6]):
        text_name, _ = _S13_COMPONENT_SHAPES[i]
        sh = _find_shape_by_name(slide, text_name)
        if sh is not None:
            _set_shape_text(sh, f"{label} ({pct}%)")
    for text_name, marker_name in _S13_COMPONENT_SHAPES[len(slices):]:
        _remove_shape(slide, text_name)
        _remove_shape(slide, marker_name)

    # Trim the background marker strip (Rounded Rectangle 18) to span only
    # the kept components. Its original height spans all 6; when fewer
    # components are used, shrink to match the last kept marker's bottom.
    bg = _find_shape_by_name(slide, "Rounded Rectangle 18")
    if bg is not None and len(slices) < 6:
        last_marker_name = _S13_COMPONENT_SHAPES[len(slices) - 1][1]
        last_marker = _find_shape_by_name(slide, last_marker_name)
        if last_marker is not None:
            bg.height = (last_marker.top + last_marker.height) - bg.top


# ── Cleanup helpers ──────────────────────────────────────────

def strip_helper_textboxes(prs):
    """Remove template instructional textboxes (e.g. 'Link to more info...') from all slides."""
    removed = 0
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if hasattr(shape, "text_frame"):
                if "link to more" in shape.text_frame.text.lower():
                    shape._element.getparent().remove(shape._element)
                    removed += 1
    return removed


def set_slide_view(prs):
    """Set the file to open in normal slide view (not master view)."""
    for rel in prs.part.rels.values():
        if "viewProps" in rel.reltype:
            xml = etree.fromstring(rel.target_part.blob)
            xml.set("lastView", "sldView")
            rel.target_part._blob = etree.tostring(
                xml, xml_declaration=True, encoding="UTF-8", standalone=True
            )
            break
