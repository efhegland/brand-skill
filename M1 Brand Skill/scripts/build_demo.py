#!/usr/bin/env python3
"""Build the M1 AI Wealth Management demo deck from testing/test-01.md.

Single output: testing/M1-AI-Wealth-Management.pptx — one coherent deck that
exercises every polish rule (S1 cover, S2 agenda, S4 section breaks, S6
quote, S7/S8 content with auto-imagery routing to S9/S10/S11, **S12 bar
chart**, **S13 pie chart**, **S14 timeline graphic**, **S15 target graphic**,
**S16 takeaways graphic**, **S17 table**, S18 closing). Use it as the
canonical visual reference for what the skill produces.

Pre-analysis pass — each subsection routed per the Graphics vs Charts rule:

| Subsection | Route to | Reason |
|---|---|---|
| II.A relationship model       | S16 takeaways graphic | 3 items with takeaways |
| II.B three pillars            | S16 takeaways graphic | Definitive 3 pillars |
| II.D wealth journey (NEW)     | S15 target graphic    | 4 nested concepts that build on each other |
| III.B three stages            | S14 timeline graphic  | 3 SEQUENTIAL stages |
| III.C client quote            | S6 quote              | Actual quote |
| IV.A growth since launch      | S12 bar chart         | Per-period dollar values |
| IV.B asset distribution       | S13 pie chart         | % adds to 100% |
| IV.C how M1 compares          | S17 table             | Matrix comparison |
| V.B rates                     | S16 takeaways graphic | 3 rate products with details |
| V.C what others say           | S16 takeaways graphic | 3 proof points |
| (everything else)             | S7/S8 content         | Plain text content |
"""
import os
import shutil
import sys

BASE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(BASE)
sys.path.insert(0, BASE)

from pptx import Presentation
from pptx_helpers import (
    duplicate_slide, delete_slide, fill_agenda, grey_box_slides,
    imagery_warnings, reset_image_tracking, routing_warnings,
    set_content_title, set_cover_title, set_content_series, set_quote,
    set_section_break, set_slide_view, strip_helper_textboxes,
    set_s12_bar, set_s13_pie,
    set_s14_timeline, set_s15_target, set_s16_takeaways,
)
from render_for_review import render_cover_to_png, render_slide_to_png

TMPL = os.path.join(ROOT, "assets", "M1-Presentation-Template_2026.pptx")
OUT  = os.path.join(ROOT, "testing", "M1-AI-Wealth-Management.pptx")

# Template slide indices (18-slide template)
S1_COVER, S2_AGENDA, S4_DARK_SB = 0, 1, 3
S6_QUOTE                        = 5
S12_BAR, S13_PIE                = 11, 12
S14_TIMELINE, S15_TARGET, S16_TAKEAWAYS = 13, 14, 15
S17_TABLE, S18_CLOSING          = 16, 17
N_TEMPLATE_SLIDES               = 18

shutil.copy(TMPL, OUT)
prs = Presentation(OUT)

# Fresh build — clear cross-deck image-tracking state so each build starts
# with no images marked as used.
reset_image_tracking()


def section(title):
    """S4 dark section break (no supporting text)."""
    s = duplicate_slide(prs, S4_DARK_SB)
    set_section_break(s, title)


# ── 1. Cover ───────────────────────────────────────────────────
s = duplicate_slide(prs, S1_COVER)
set_cover_title(s, "M1: AI Wealth Management")

# ── 2. Agenda ──────────────────────────────────────────────────
s = duplicate_slide(prs, S2_AGENDA)
ph = next(p for p in s.placeholders if p.placeholder_format.idx == 1)
fill_agenda(ph, [
    "Why personal finance is broken",
    "M1 is the third option",
    "M1 Advisor",
    "Platform traction",
    "The M1 advantage",
])

# ── I. Why personal finance is broken ──────────────────────────
section("Why personal finance is broken")

# I.A → plain content (dense; cadence picks the best-match image elsewhere)
set_content_series(prs, [
    ("The system wasn't designed for you", [
        (0, "Traditional wealth management was built for the ultra-wealthy"),
        (0, "Everyone else inherited a fragmented patchwork of products"),
        (2, "A brokerage here, a savings account there, a loan from somewhere else"),
        (2, "No one manages the whole picture"),
        (2, "You end up bridging the gaps yourself — or they go unbridged"),
    ]),
])

# I.B → S16 takeaways (title number-keyword "two" — graphic-first routing)
# 2 parallel discrete topics; the 3rd row is auto-removed by the helper
s = duplicate_slide(prs, S16_TAKEAWAYS)
set_s16_takeaways(s, "The two gaps you're left to fill", [
    ("Expertise gap", "You're expected to be the expert on every part of your financial life"),
    ("Execution gap", [
        ("Rebalancing",  "Portfolios kept aligned manually"),
        ("Moving money", "Transfers between accounts by hand"),
        ("Tracking",     "Performance pieced together across disconnected apps"),
    ]),
])

# I.C → plain content; routed through `set_content_series` so the cadence
# image-promotion picks it up. Title number-keyword "Two" hinted graphic, but
# I.B already took the graphic slot — avoiding two consecutive S16 slides.
# The screenshot annotation on this slide explicitly highlighted the
# "lots of empty space" → image hint, so image variant is the right fallback.
set_content_series(prs, [
    ("Two models — only one has worked", [
        (0, "Transactional model: fragmented products, total self-reliance"),
        (0, "Relationship model: integrated platform, aligned expertise"),
        (2, "Historically required $10M or more in investable assets"),
        (2, "A private banker handles strategy and execution"),
        (2, "Not available to most Americans"),
    ]),
])

# ── II. M1 is the third option ─────────────────────────────────
# Note on slide-type variety: II.A and II.B both fit S16 takeaways
# structurally, but two S16 in a row would feel repetitive. II.A is the
# descriptive setup (3 capabilities) — render as plain content so II.B can
# carry the definitive "3 pillars" S16 treatment with maximum impact.
section("M1 is the third option")

# II.A → plain content (avoid consecutive S16 with II.B below)
set_content_series(prs, [
    ("The relationship model — without the relationship cost", [
        (0, "Intelligent software replaces expensive human expertise"),
        (0, "An integrated platform replaces fragmented point solutions"),
        (0, "You get comprehensive wealth management — once reserved for clients of private banks"),
    ]),
])

# II.B → S16 takeaways graphic (definitive 3 pillars)
s = duplicate_slide(prs, S16_TAKEAWAYS)
set_s16_takeaways(s, "The three pillars of your financial life", [
    ("Cash flow", "Earn more on your money while staying liquid"),
    ("Balance sheet", "Invest for the long term, borrow when it makes sense"),
    ("Stewardship", "Protect what you've built and optimize how it grows"),
])

# II.C → plain content (2 items with nested bullets)
set_content_series(prs, [
    ("How the ecosystem works together", [
        (0, "Each product is strong on its own"),
        (0, "Together, they create something none could deliver alone"),
        (2, "Money moves intelligently between accounts based on your rules"),
        (2, "Every recommendation draws from your complete financial picture"),
        (2, "Automation handles the busy work — you set the strategy"),
    ]),
])

# II.D → S15 target graphic (4 nested concepts, innermost → outermost)
s = duplicate_slide(prs, S15_TARGET)
set_s15_target(s, "The wealth journey, layer by layer", [
    ("Foundation", "Emergency cash that earns more than it sits idle"),  # Core
    ("Growth",     "Invested portfolios that compound across decades"),   # Need
    ("Leverage",   "Borrowing power without forced selling"),             # Utility
    ("Legacy",     "Stewardship that protects what you've built"),         # Desire
])

# ── III. M1 Advisor ────────────────────────────────────────────
section("M1 Advisor")

# III.A → plain content (3 items — descriptive, no clear "guiding number" structure)
set_content_series(prs, [
    ("Your AI Wealth Management partner", [
        (0, "Understands your goals and tracks your progress against them"),
        (0, "Surfaces the right insight at the right moment"),
        (0, "Prepares specific recommendations — you retain control"),
    ]),
])

# III.B → S14 timeline graphic (3 SEQUENTIAL stages)
s = duplicate_slide(prs, S14_TIMELINE)
set_s14_timeline(
    s,
    "The M1 Advisor maturity model",
    [
        ("Inform",  "Answers your questions and builds your financial literacy"),
        ("Analyze", "Turns raw data into signals you can act on"),
        ("Guide",   "Prepares next-step recommendations for your approval"),
    ],
    takeaway="Each stage compounds on the last — M1 Advisor gets sharper the longer you use it.",
)

# III.C → quote
qs = duplicate_slide(prs, S6_QUOTE)
set_quote(qs, "Managing money used to feel overwhelming. With M1, it just runs.", "M1 client")

# ── IV. Platform traction ──────────────────────────────────────
section("Platform traction")

# IV.A → S12 bar chart (per-period dollar values from outline)
bs = duplicate_slide(prs, S12_BAR)
set_s12_bar(bs, "Growth since launch", [
    ("Q3 2019", 0),    # launch
    ("Q2 2021", 5),
    ("Q1 2023", 8),
    ("Q4 2025", 12),
], series_name="Client assets ($B)")

# IV.B → S13 pie (% adds to 100%)
ps = duplicate_slide(prs, S13_PIE)
set_s13_pie(ps, "Where client assets live today", [
    ("Invest accounts",       68),
    ("High-Yield Cash",       18),
    ("Borrow (margin/loans)", 14),
])

# IV.C → SKIPPED. The S17 table requires a comparison matrix (rows × cols
# of M1 vs. traditional brokerage vs. robo-advisor features) and the
# outline only supplies the title intent. Per the no-empty-slides
# directive, we don't ship a title-only or empty-body placeholder.
# Surface the gap via a routing warning instead.
routing_warnings.append(
    "IV.C 'How M1 compares' SKIPPED — outline lacks the comparison matrix "
    "cell data. Add rows × columns (M1 / traditional / robo-advisor) to "
    "the outline to enable the S17 table treatment. The slide is omitted "
    "from the deck entirely until that data exists."
)

# ── V. The M1 advantage ────────────────────────────────────────
section("The M1 advantage")

# V.A → plain content (2 items with nested)
set_content_series(prs, [
    ("Automation that works on your terms", [
        (0, "Dynamic Rebalancing keeps portfolios aligned without forcing you to sell"),
        (0, "Trade windows protect long-term investors from reactive decisions"),
        (2, "Enforces discipline by design"),
        (2, "Prevents emotional trading during market swings"),
    ]),
])

# V.B → S16 takeaways graphic (mixed: 1-proof, 2-proof, 4-proof rows
# demonstrating the flex layout)
s = duplicate_slide(prs, S16_TAKEAWAYS)
set_s16_takeaways(s, "Rates built for wealth-builders", [
    # Topic 1: 1 wide proof
    ("High-Yield Cash", "APY that leads national averages — no minimums, no monthly fees"),
    # Topic 2: 2-column proof row
    ("Margin Loan", [
        ("Rates",       "Competitive with institutional lenders"),
        ("Flexibility", "Borrow against your portfolio without selling"),
    ]),
    # Topic 3: 4-column proof row (the template's max for a single-row layout)
    ("Personal Loan", [
        ("Fixed APR",     "Starting at 7.99% with no markup over the term"),
        ("No fees",       "No origination, no prepayment, no late fees"),
        ("Funding speed", "Approved in minutes, funded same business day"),
        ("Term options",  "36 / 48 / 60 month terms to fit your cash flow"),
    ]),
])

# V.C → plain content (avoid consecutive S16 with V.B above; 3 proof points
# work fine as a bulleted list)
set_content_series(prs, [
    ("What others say about M1", [
        (0, "Investopedia named M1 the best brokerage for sophisticated investors (2024)"),
        (0, "$12B+ in client assets managed for hundreds of thousands of clients"),
        (0, '"You\'re serious about your money. So are we."'),
    ]),
])

# ── Closing ────────────────────────────────────────────────────
duplicate_slide(prs, S18_CLOSING)

# ── Delete originals, cleanup, save ────────────────────────────
for i in range(N_TEMPLATE_SLIDES - 1, -1, -1):
    delete_slide(prs, i)
strip_helper_textboxes(prs)
set_slide_view(prs)
prs.save(OUT)

# ── Summary ────────────────────────────────────────────────────
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
if imagery_warnings:
    print(f"\nImagery warnings ({len(imagery_warnings)}):")
    for w in imagery_warnings:
        print(f"  ⚠ {w}")
if routing_warnings:
    print(f"\nRouting warnings ({len(routing_warnings)}):")
    for w in routing_warnings:
        print(f"  ⚠ {w}")
if grey_box_slides:
    print(f"\nSlides needing user-supplied images ({len(grey_box_slides)}):")
    for title in grey_box_slides:
        print(f"  □ {title}")
    print("  → OneDrive image catalog was unavailable for this build. "
          "Open the deck and drop an image into each placeholder above "
          "before sharing.")

# ── Visual verification — render text-fitting and visual slides ──
print()
to_render = []
final_prs = Presentation(OUT)
for i, slide in enumerate(final_prs.slides):
    layout = slide.slide_layout.name
    if i == 0 or layout in ("SectionBreak", "1_Title Slide", "1_SectionBreak"):
        to_render.append((i, f"slide{i+1} ({layout})"))
    elif layout == "Content Text " and i == 7:
        to_render.append((i, f"slide{i+1} (multi-line content title)"))

for idx, label in to_render:
    try:
        if idx == 0:
            png = render_cover_to_png(OUT)
        else:
            png = render_slide_to_png(OUT, idx)
        print(f"Verify {label}: {png}")
    except Exception as e:
        print(f"⚠ Render failed for {label}: {e}")
